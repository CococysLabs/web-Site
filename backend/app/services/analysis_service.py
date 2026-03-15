"""
Servicio de análisis completo de documentos con Gemini.
Soporta: PDF, DOCX, PPTX, XLSX, TXT y archivos nativos de Google Workspace
(Docs → DOCX, Sheets → XLSX, Slides → PPTX, exportados automáticamente por drive_service).
"""
from typing import Dict, List, Optional, Any
import io
import json
import re
import math
from datetime import datetime

try:
    import google.generativeai as genai
    GENAI_AVAILABLE = True
except ImportError:
    GENAI_AVAILABLE = False

from app.config import settings

GEMINI_MODEL = "gemini-2.0-flash"

# Formatos soportados para análisis (incluyendo Google Workspace nativos)
SUPPORTED_MIMES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",   # DOCX
    "application/vnd.openxmlformats-officedocument.presentationml.presentation", # PPTX
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",         # XLSX
    "text/plain",
    # Google Workspace (drive_service los exporta automáticamente)
    "application/vnd.google-apps.document",
    "application/vnd.google-apps.spreadsheet",
    "application/vnd.google-apps.presentation",
}

# Para el análisis con IA, usar este número máximo de caracteres
AI_ANALYSIS_CHARS = 14_000
# Palabras por minuto para calcular tiempo de lectura
WPM = 200


class DocumentAnalysisService:
    """
    Analiza documentos de todos los formatos con Gemini AI.
    Produce: existencia/metadatos, estructura, contexto y evaluación de calidad.
    """

    def __init__(self):
        self.model = None
        self.enabled = False
        if GENAI_AVAILABLE:
            api_key = getattr(settings, 'GEMINI_API_KEY', None)
            if api_key and api_key not in ('', 'kjkj'):
                genai.configure(api_key=api_key)
                self.model = genai.GenerativeModel(GEMINI_MODEL)
                self.enabled = True
                print(f"✅ Gemini ({GEMINI_MODEL}) listo para análisis de documentos")
            else:
                print("⚠️  GEMINI_API_KEY no configurada — modo básico")
        else:
            print("⚠️  google-generativeai no disponible — modo básico")

    # ──────────────────────────────────────────────────────────────────────────
    # Extracción de texto + metadatos por formato
    # ──────────────────────────────────────────────────────────────────────────

    def extract_text(self, content: bytes, mime_type: str) -> str:
        """Extrae texto del documento según su tipo MIME."""
        if "pdf" in mime_type:
            return self._extract_pdf(content)["text"]
        elif "wordprocessingml" in mime_type:
            return self._extract_docx(content)["text"]
        elif "presentationml" in mime_type:
            return self._extract_pptx(content)["text"]
        elif "spreadsheetml" in mime_type:
            return self._extract_xlsx(content)["text"]
        elif "text/plain" in mime_type:
            return self._extract_txt(content)
        return ""

    def _extract_pdf(self, content: bytes) -> Dict[str, Any]:
        """Extrae texto y metadatos de un PDF."""
        from PyPDF2 import PdfReader
        try:
            reader = PdfReader(io.BytesIO(content))
            pages_text = []
            total_chars = 0
            has_images = False

            for i, page in enumerate(reader.pages, 1):
                page_text = page.extract_text() or ""
                total_chars += len(page_text)
                if page_text.strip():
                    pages_text.append(f"=== Página {i} ===\n{page_text}")

                # Detectar imágenes (best-effort)
                if not has_images:
                    try:
                        resources = page.get("/Resources")
                        if resources:
                            xobj = resources.get("/XObject")
                            if xobj:
                                xobj_obj = xobj.get_object() if hasattr(xobj, "get_object") else xobj
                                for k in xobj_obj:
                                    sub = xobj_obj[k]
                                    if hasattr(sub, "get_object"):
                                        sub = sub.get_object()
                                    if sub.get("/Subtype") == "/Image":
                                        has_images = True
                                        break
                    except Exception:
                        pass

            metadata = {}
            try:
                if reader.metadata:
                    metadata = {k.lstrip("/"): v for k, v in reader.metadata.items()}
            except Exception:
                pass

            full_text = "\n".join(pages_text)
            return {
                "text": full_text,
                "pages": len(reader.pages),
                "total_chars": total_chars,
                "has_images": has_images,
                "metadata": metadata,
                "readable": total_chars > 50,
                "file_size_kb": round(len(content) / 1024, 2),
            }
        except Exception as e:
            return {"text": "", "pages": 0, "total_chars": 0, "has_images": False,
                    "metadata": {}, "readable": False, "file_size_kb": round(len(content)/1024, 2), "error": str(e)}

    def _extract_docx(self, content: bytes) -> Dict[str, Any]:
        """Extrae texto y metadatos de un DOCX."""
        try:
            from docx import Document as DocxDoc
            doc = DocxDoc(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # También incluir texto de tablas
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            paragraphs.append(cell.text.strip())
            full_text = "\n".join(paragraphs)
            # Contar imágenes embebidas
            has_images = len(doc.inline_shapes) > 0
            return {
                "text": full_text,
                "pages": None,       # DOCX no tiene páginas fijas
                "paragraph_count": len(paragraphs),
                "table_count": len(doc.tables),
                "total_chars": len(full_text),
                "has_images": has_images,
                "readable": len(full_text) > 50,
                "file_size_kb": round(len(content) / 1024, 2),
            }
        except Exception as e:
            return {"text": "", "paragraph_count": 0, "table_count": 0, "total_chars": 0,
                    "has_images": False, "readable": False, "file_size_kb": round(len(content)/1024, 2), "error": str(e)}

    def _extract_pptx(self, content: bytes) -> Dict[str, Any]:
        """Extrae texto y metadatos de un PPTX."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            slides_text = []
            has_images = False
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        has_images = True
                if parts:
                    slides_text.append(f"=== Diapositiva {i} ===\n" + "\n".join(parts))
                # También notas del presentador
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slides_text.append(f"[Notas diap. {i}]: {notes}")
            full_text = "\n\n".join(slides_text)
            return {
                "text": full_text,
                "slides": len(prs.slides),
                "pages": len(prs.slides),
                "total_chars": len(full_text),
                "has_images": has_images,
                "readable": len(full_text) > 50,
                "file_size_kb": round(len(content) / 1024, 2),
            }
        except Exception as e:
            return {"text": "", "slides": 0, "pages": 0, "total_chars": 0,
                    "has_images": False, "readable": False, "file_size_kb": round(len(content)/1024, 2), "error": str(e)}

    def _extract_xlsx(self, content: bytes) -> Dict[str, Any]:
        """Extrae texto y metadatos de un XLSX."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            parts = []
            total_rows = 0
            sheet_summaries = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows_text = []
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    # Tomar solo las primeras 200 filas por hoja para no saturar
                    if row_count >= 200:
                        break
                    row_vals = [str(v) for v in row if v is not None and str(v).strip()]
                    if row_vals:
                        rows_text.append("\t".join(row_vals))
                        row_count += 1
                total_rows += row_count
                sheet_summaries.append(f"{sheet_name} ({row_count} filas)")
                if rows_text:
                    parts.append(f"=== Hoja: {sheet_name} ===\n" + "\n".join(rows_text))
            full_text = "\n\n".join(parts)
            return {
                "text": full_text,
                "sheets": wb.sheetnames,
                "sheet_count": len(wb.sheetnames),
                "total_rows": total_rows,
                "sheet_summaries": sheet_summaries,
                "pages": len(wb.sheetnames),
                "total_chars": len(full_text),
                "has_images": False,
                "readable": len(full_text) > 10,
                "file_size_kb": round(len(content) / 1024, 2),
            }
        except Exception as e:
            return {"text": "", "sheets": [], "sheet_count": 0, "total_rows": 0,
                    "has_images": False, "readable": False, "file_size_kb": round(len(content)/1024, 2), "error": str(e)}

    def _extract_txt(self, content: bytes) -> str:
        """Extrae texto de un archivo de texto plano."""
        for enc in ("utf-8", "latin-1", "utf-16"):
            try:
                return content.decode(enc)
            except Exception:
                continue
        return content.decode("utf-8", errors="replace")

    def _extract_all(self, content: bytes, mime_type: str) -> Dict[str, Any]:
        """
        Extrae texto + metadatos según el mime_type.
        Retorna dict con: text, word_count, reading_time_min, file_size_kb, pages,
        has_images, readable, y claves específicas del formato.
        """
        if "pdf" in mime_type:
            data = self._extract_pdf(content)
        elif "wordprocessingml" in mime_type:
            data = self._extract_docx(content)
        elif "presentationml" in mime_type:
            data = self._extract_pptx(content)
        elif "spreadsheetml" in mime_type:
            data = self._extract_xlsx(content)
        elif "text/plain" in mime_type:
            text = self._extract_txt(content)
            data = {
                "text": text,
                "pages": None,
                "total_chars": len(text),
                "has_images": False,
                "readable": len(text) > 20,
                "file_size_kb": round(len(content) / 1024, 2),
            }
        else:
            data = {"text": "", "pages": None, "total_chars": 0,
                    "has_images": False, "readable": False,
                    "file_size_kb": round(len(content) / 1024, 2)}

        text = data.get("text", "")
        word_count = len(text.split())
        reading_time_min = max(1, math.ceil(word_count / WPM))

        return {
            **data,
            "word_count": word_count,
            "reading_time_min": reading_time_min,
        }

    # Legacy aliases
    def extract_text_from_pdf(self, pdf_content: bytes) -> str:
        return self._extract_pdf(pdf_content)["text"]

    def extract_text(self, content: bytes, mime_type: str) -> str:
        return self._extract_all(content, mime_type).get("text", "")

    # ──────────────────────────────────────────────────────────────────────────
    # Análisis principal: análisis completo con un solo llamado a Gemini
    # ──────────────────────────────────────────────────────────────────────────

    def analyze_complete(self, content: bytes, mime_type: str = "application/pdf") -> Dict[str, Any]:
        """
        Análisis completo: existencia + estructura + contexto + calidad académica.
        Usa un solo llamado a Gemini (más coherente y eficiente).
        """
        extracted = self._extract_all(content, mime_type)
        text = extracted.get("text", "")
        word_count = extracted.get("word_count", 0)

        # Construir objeto de existencia
        existence = {
            "exists": True,
            "readable": extracted.get("readable", False),
            "pages": extracted.get("pages"),
            "slides": extracted.get("slides"),
            "sheet_count": extracted.get("sheet_count"),
            "sheets": extracted.get("sheets"),
            "file_size_kb": extracted.get("file_size_kb", 0),
            "has_content": word_count > 20,
            "total_characters": extracted.get("total_chars", 0),
            "has_images": extracted.get("has_images", False),
            "metadata": extracted.get("metadata", {}),
            "word_count": word_count,
            "reading_time_min": extracted.get("reading_time_min", 0),
            "paragraph_count": extracted.get("paragraph_count"),
            "table_count": extracted.get("table_count"),
            "sheet_summaries": extracted.get("sheet_summaries"),
            "error": extracted.get("error"),
        }

        if not text or word_count < 5:
            return {
                "existence": existence,
                "structure": self._empty_structure(),
                "context": self._empty_context(),
                "quality": self._empty_quality(),
                "analyzed_at": datetime.utcnow().isoformat(),
                "gemini_enabled": self.enabled,
                "mime_type": mime_type,
            }

        # Determinar tipo de documento para el prompt
        doc_type_hint = self._mime_to_hint(mime_type)

        if self.enabled:
            ai_result = self._comprehensive_ai_analysis(text, word_count, doc_type_hint)
        else:
            ai_result = self._comprehensive_basic_analysis(text)

        return {
            "existence": existence,
            "structure": ai_result.get("structure", self._empty_structure()),
            "context":   ai_result.get("context",   self._empty_context()),
            "quality":   ai_result.get("quality",   self._empty_quality()),
            "analyzed_at": datetime.utcnow().isoformat(),
            "gemini_enabled": self.enabled,
            "mime_type": mime_type,
        }

    def _mime_to_hint(self, mime_type: str) -> str:
        if "pdf" in mime_type:
            return "PDF"
        elif "wordprocessingml" in mime_type or "google-apps.document" in mime_type:
            return "documento Word/Google Doc"
        elif "presentationml" in mime_type or "google-apps.presentation" in mime_type:
            return "presentación PowerPoint/Google Slides"
        elif "spreadsheetml" in mime_type or "google-apps.spreadsheet" in mime_type:
            return "hoja de cálculo Excel/Google Sheets"
        elif "text/plain" in mime_type:
            return "archivo de texto"
        return "documento"

    # ──────────────────────────────────────────────────────────────────────────
    # Análisis con IA — prompt único comprehensivo
    # ──────────────────────────────────────────────────────────────────────────

    def _comprehensive_ai_analysis(self, text: str, word_count: int, doc_type_hint: str) -> Dict[str, Any]:
        """
        Un solo llamado a Gemini con prompt exhaustivo.
        Analiza estructura, contexto y calidad académica en una sola respuesta.
        """
        # Truncar para no exceder límites del modelo, pero dar suficiente contexto
        analysis_text = text[:AI_ANALYSIS_CHARS]
        truncated = len(text) > AI_ANALYSIS_CHARS

        prompt = f"""Eres un evaluador académico universitario experto. Analiza exhaustivamente el siguiente {doc_type_hint} ({word_count} palabras{"—se muestra la primera parte" if truncated else ""}).

CONTENIDO DEL DOCUMENTO:
{analysis_text}
{"[... el documento continúa ...]" if truncated else ""}

Responde ÚNICAMENTE con JSON válido sin texto adicional ni markdown:
{{
  "structure": {{
    "has_table_of_contents": true|false,
    "sections": ["Sección o título detectado", ...],
    "has_bibliography": true|false,
    "has_tables": true|false,
    "has_images": true|false,
    "total_sections": 0,
    "analysis_method": "ai"
  }},
  "context": {{
    "summary": "Resumen conciso de 2-3 oraciones sobre el contenido",
    "main_topics": ["Tema principal 1", "Tema 2", ...],
    "language": "español|inglés|otro",
    "academic_level": "pregrado|posgrado|técnico|profesional|no determinado",
    "document_type": "syllabus|tesis|proyecto|informe|presentación|material de clase|hoja de datos|guía|otro",
    "keywords": ["palabra clave 1", "kw2", ...],
    "word_count": {word_count},
    "analysis_method": "ai"
  }},
  "quality": {{
    "score": 0,
    "level": "excelente|bueno|regular|bajo|sin contenido",
    "strengths": ["Fortaleza detectada 1", ...],
    "weaknesses": ["Área de mejora 1", ...],
    "recommendations": ["Recomendación concreta 1", ...],
    "completeness": "completo|parcial|incompleto",
    "analysis_method": "ai"
  }}
}}

Para el campo "quality.score" usa una escala 0-100:
- 90-100: Documento excelente, completo y bien estructurado
- 70-89: Buen documento con algunos aspectos a mejorar
- 50-69: Documento regular, faltan elementos importantes
- 30-49: Documento incompleto o con problemas serios
- 0-29: Documento muy deficiente o con mínimo contenido útil
"""

        try:
            response = self.model.generate_content(prompt)
            raw = response.text.strip()
            raw = re.sub(r'^```json?\s*', '', raw)
            raw = re.sub(r'\s*```$', '', raw)
            parsed = json.loads(raw)
            return parsed
        except json.JSONDecodeError:
            print(f"  ⚠️  JSON inválido de Gemini en análisis completo, usando básico")
            return self._comprehensive_basic_analysis(text)
        except Exception as e:
            print(f"  ⚠️  Error en Gemini análisis completo: {e}")
            return self._comprehensive_basic_analysis(text)

    # ──────────────────────────────────────────────────────────────────────────
    # Análisis básico sin IA (fallback)
    # ──────────────────────────────────────────────────────────────────────────

    def _comprehensive_basic_analysis(self, text: str) -> Dict[str, Any]:
        """Análisis heurístico sin IA."""
        lines = text.split('\n')
        text_lower = text.lower()
        word_count = len(text.split())

        # Detectar secciones
        sections = []
        for i, line in enumerate(lines):
            clean = line.strip()
            if not clean or len(clean) > 120:
                continue
            if clean.isupper() and len(clean) > 3:
                sections.append(clean)
            elif re.match(r'^\d+\.?\s+[A-ZÁÉÍÓÚ]', clean):
                sections.append(clean)

        has_toc = any('contenido' in s.lower() or 'índice' in s.lower() for s in sections)
        has_bib = any(kw in text_lower for kw in ('bibliografía', 'referencias', 'bibliography'))
        has_tbl = any(kw in text_lower for kw in ('tabla', 'cuadro', 'table'))
        has_img = any(kw in text_lower for kw in ('figura', 'imagen', 'fig.', 'gráfico'))

        # Detectar idioma
        es_count = sum(text_lower.count(w) for w in ['el', 'la', 'de', 'que', 'en', 'y', 'los'])
        en_count = sum(text_lower.count(w) for w in ['the', 'of', 'and', 'to', 'a', 'in', 'is'])
        language = 'español' if es_count > en_count else 'inglés'

        # Tipo de documento
        doc_type = 'documento'
        for keyword, dtype in [('syllabus', 'syllabus'), ('tesis', 'tesis'),
                                ('proyecto', 'proyecto'), ('informe', 'informe'),
                                ('presentación', 'presentación')]:
            if keyword in text_lower:
                doc_type = dtype
                break

        # Calidad heurística
        score = 0
        if word_count > 500:  score += 30
        elif word_count > 100: score += 15
        if sections:          score += 20
        if has_bib:           score += 15
        if has_toc:           score += 10
        if has_tbl:           score += 10
        if has_img:           score += 10
        score = min(score, 85)  # Max 85 sin IA

        words_sample = text.split()[:80]
        summary = ' '.join(words_sample) + ('...' if len(text.split()) > 80 else '')

        return {
            "structure": {
                "has_table_of_contents": has_toc,
                "sections": sections[:20],
                "has_bibliography": has_bib,
                "has_tables": has_tbl,
                "has_images": has_img,
                "total_sections": len(sections),
                "analysis_method": "basic",
            },
            "context": {
                "summary": summary,
                "main_topics": [],
                "language": language,
                "academic_level": "no determinado",
                "document_type": doc_type,
                "keywords": [],
                "word_count": word_count,
                "analysis_method": "basic",
            },
            "quality": {
                "score": score,
                "level": "bueno" if score >= 70 else "regular" if score >= 40 else "bajo",
                "strengths": (["Contiene bibliografía"] if has_bib else []) +
                             (["Estructura con secciones"] if sections else []),
                "weaknesses": (["Sin bibliografía detectada"] if not has_bib else []) +
                              (["Sin tabla de contenidos"] if not has_toc else []),
                "recommendations": ["Usar análisis con IA para evaluación más detallada"],
                "completeness": "completo" if score >= 70 else "parcial" if score >= 40 else "incompleto",
                "analysis_method": "basic",
            }
        }

    # ──────────────────────────────────────────────────────────────────────────
    # Métodos públicos individuales (backwards compat con analysis.py route)
    # ──────────────────────────────────────────────────────────────────────────

    def analyze_existence(self, content: bytes, mime_type: str = "application/pdf") -> Dict[str, Any]:
        extracted = self._extract_all(content, mime_type)
        return {
            "exists": True,
            "readable": extracted.get("readable", False),
            "pages": extracted.get("pages"),
            "file_size_kb": extracted.get("file_size_kb", 0),
            "has_content": extracted.get("word_count", 0) > 20,
            "total_characters": extracted.get("total_chars", 0),
            "has_images": extracted.get("has_images", False),
            "metadata": extracted.get("metadata", {}),
            "word_count": extracted.get("word_count", 0),
            "error": extracted.get("error"),
        }

    def analyze_structure(self, content: bytes, mime_type: str = "application/pdf") -> Dict[str, Any]:
        extracted = self._extract_all(content, mime_type)
        text = extracted.get("text", "")
        if not text:
            return self._empty_structure()
        if self.enabled:
            result = self._comprehensive_ai_analysis(text, extracted.get("word_count", 0), self._mime_to_hint(mime_type))
            return result.get("structure", self._empty_structure())
        return self._comprehensive_basic_analysis(text).get("structure", self._empty_structure())

    def analyze_context(self, content: bytes, mime_type: str = "application/pdf") -> Dict[str, Any]:
        extracted = self._extract_all(content, mime_type)
        text = extracted.get("text", "")
        if not text:
            return self._empty_context()
        if self.enabled:
            result = self._comprehensive_ai_analysis(text, extracted.get("word_count", 0), self._mime_to_hint(mime_type))
            return result.get("context", self._empty_context())
        return self._comprehensive_basic_analysis(text).get("context", self._empty_context())

    # ──────────────────────────────────────────────────────────────────────────
    # Helpers vacíos
    # ──────────────────────────────────────────────────────────────────────────

    def _empty_structure(self) -> Dict[str, Any]:
        return {
            "has_table_of_contents": False, "sections": [], "has_bibliography": False,
            "has_tables": False, "has_images": False, "total_sections": 0, "analysis_method": "none"
        }

    def _empty_context(self) -> Dict[str, Any]:
        return {
            "summary": "No se pudo analizar el contenido", "main_topics": [],
            "language": "desconocido", "academic_level": "no determinado",
            "document_type": "desconocido", "keywords": [], "word_count": 0, "analysis_method": "none"
        }

    def _empty_quality(self) -> Dict[str, Any]:
        return {
            "score": 0, "level": "sin contenido", "strengths": [], "weaknesses": [],
            "recommendations": [], "completeness": "incompleto", "analysis_method": "none"
        }


# Singleton
analysis_service = DocumentAnalysisService()
