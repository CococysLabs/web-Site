"""
Servicio para validar el CONTENIDO de documentos contra los requisitos
de la hoja "Matriz observaciones" del Excel de la carpeta de curso.

Flujo:
  1. Deriva la sección desde el nombre de la carpeta Semana_X
  2. Lee la 2da hoja del Excel con openpyxl
  3. Filtra los requisitos de esa sección donde Aplica == "Si"
  4. Descarga todos los documentos (PDF/DOCX/PPTX) de la carpeta
  5. Extrae texto de cada documento
  6. Usa Gemini AI para verificar si cada sub-sección está implícitamente cubierta
  7. Escribe resultados ("Presente", "Observaciones") de vuelta al Excel en Drive
"""
import io
import re
import json
import unicodedata
import time
from typing import Dict, List, Optional, Any, Tuple

import openpyxl
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

try:
    import google.generativeai as genai
    GENAI_AVAILABLE = True
except ImportError:
    GENAI_AVAILABLE = False

from app.config import settings
from app.services.drive_service import drive_service


# ─── Constantes ──────────────────────────────────────────────────────────────

MATRIX_FILE_PREFIX = "Matriz observaciones"
GEMINI_MODEL       = "gemini-2.0-flash"
MAX_CHUNK_CHARS    = 25_000
CHUNK_OVERLAP      = 500
EXCEL_MIME_TYPE    = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

SUPPORTED_MIMES = {
    # Formatos binarios directos
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    # Google Workspace nativos — drive_service los exporta automáticamente
    "application/vnd.google-apps.presentation",   # → PPTX
    "application/vnd.google-apps.document",        # → DOCX
    "application/vnd.google-apps.spreadsheet",     # → XLSX (no se usa para extracción)
}

# Vocabulario de tipos de documento que actúan como ENCABEZADOS de grupo en el Excel.
# Solo nombres que representan el TIPO de archivo a buscar en Drive, no parámetros
# de contenido. Cuestionario/Actividad/Tarea/etc. son parámetros DENTRO de un tipo.
#
# Regla de inclusión: ¿corresponde a un archivo separado en Drive?
#   ✓ Presentación → Presentacion_S2.pptx
#   ✓ Lectura      → Lectura_S2.pdf
#   ✓ Video        → Video_S2.mp4 (o cuestionario PDF asociado al video)
#   ✗ Cuestionario → es un parámetro/requisito DENTRO del Video, no un archivo propio
#   ✗ Actividad    → ídem
DOCUMENT_TYPE_HEADERS: set = {
    'presentacion', 'presentación',
    'diapositivas', 'slides',
    'lectura', 'reading', 'documento',
    'video',
    'guia', 'guía',
    'informe',
    'reporte',
}

# Carpetas de laboratorio → nombre canónico de hoja en el Excel
LAB_FOLDER_TYPES: Dict[str, str] = {
    'proyectos': 'Proyectos',
    'practicas': 'Practicas',
    'tareas':    'Tareas',
}

# Descripción de cada tipo de carpeta lab para el prompt de IA
LAB_FOLDER_DESCRIPTIONS: Dict[str, Dict[str, str]] = {
    'proyectos': {
        'tipo':        'Proyecto Académico',
        'descripcion': 'Trabajo integrador donde el estudiante desarrolla e implementa una solución completa a un problema real.',
        'enfoque':     'Integración de conocimientos, diseño de soluciones, implementación práctica y documentación técnica.',
        'complejidad': 'Alta — requiere planificación, implementación y entrega de un producto funcional.',
    },
    'practicas': {
        'tipo':        'Práctica de Laboratorio',
        'descripcion': 'Ejercicio guiado de aplicación de técnicas o herramientas en un entorno controlado.',
        'enfoque':     'Aplicación directa de conceptos teóricos, experimentación y verificación de resultados.',
        'complejidad': 'Media — sigue pasos definidos con espacio para experimentación.',
    },
    'tareas': {
        'tipo':        'Tarea / Asignación',
        'descripcion': 'Actividad de refuerzo y evaluación de comprensión de los conceptos vistos en clase.',
        'enfoque':     'Consolidación de conocimiento, resolución de problemas y ejercicios de aplicación.',
        'complejidad': 'Básica a media — centrada en comprensión y aplicación puntual de conceptos.',
    },
}


class DocumentContentValidationService:
    """
    Valida el contenido de documentos en carpetas Semana_X contra los
    requisitos de la hoja 'Matriz observaciones' del Excel del curso.
    """

    def __init__(self):
        self.model = None
        self.enabled = False
        # Multi-key rotation: lista de todas las claves Gemini disponibles
        self._default_api_keys: List[str] = []
        self._api_keys: List[str] = []
        # Claves que ya recibieron 429 en esta sesión (se resetean al reiniciar)
        self._exhausted_keys: set = set()
        # Proveedores alternativos
        self._deepseek_key: Optional[str] = None
        self._groq_key: Optional[str] = None
        self._openrouter_key: Optional[str] = None
        self._key_source: str = "none"
        self._init_gemini()

    def _init_gemini(self):
        if not GENAI_AVAILABLE:
            print("⚠️  google-generativeai no disponible")
            return

        # ── Recolectar todas las claves Gemini ───────────────────────────────
        keys: List[str] = []

        # Clave principal (GEMINI_API_KEY)
        primary = getattr(settings, 'GEMINI_API_KEY', '') or ''
        if primary and primary not in ('', 'kjkj'):
            keys.append(primary.strip())

        # Claves adicionales (GEMINI_API_KEYS = "key1,key2,key3")
        extra = getattr(settings, 'GEMINI_API_KEYS', '') or ''
        for k in extra.split(','):
            k = k.strip()
            if k and k not in keys:
                keys.append(k)

        self._default_api_keys = list(keys)
        self._api_keys = list(keys)
        self.enabled = len(keys) > 0
        self._key_source = "env" if keys else "none"

        if self.enabled:
            # Configurar con la primera clave por defecto
            genai.configure(api_key=keys[0])
            self.model = genai.GenerativeModel(GEMINI_MODEL)
            total = len(keys)
            print(f"✅ Gemini ({GEMINI_MODEL}) listo — {total} clave{'s' if total > 1 else ''} disponible{'s' if total > 1 else ''}")
        else:
            print("⚠️  GEMINI_API_KEY no configurada — modo fallback (keyword matching)")

        # ── DeepSeek como proveedor principal ────────────────────────────────
        deepseek_key = getattr(settings, 'DEEPSEEK_API_KEY', '') or ''
        if deepseek_key:
            self._deepseek_key = deepseek_key.strip()
            print("✅ DeepSeek configurado como proveedor PRINCIPAL de IA")

        # ── Groq como proveedor de respaldo ──────────────────────────────────
        groq_key = getattr(settings, 'GROQ_API_KEY', '') or ''
        if groq_key:
            self._groq_key = groq_key.strip()
            print("✅ Groq configurado como proveedor de respaldo")

        # ── OpenRouter como proveedor de respaldo adicional ───────────────────
        openrouter_key = getattr(settings, 'OPENROUTER_API_KEY', '') or ''
        if openrouter_key:
            self._openrouter_key = openrouter_key.strip()
            print("✅ OpenRouter configurado como proveedor de respaldo adicional")

    # ──────────────────────────────────────────────────────────────────────────
    # Multi-key rotation & proveedores alternativos
    # ──────────────────────────────────────────────────────────────────────────

    @staticmethod
    def _batch_prompt(params: list, section_name: str, group_name: str, doc_snippet: str,
                      max_chars: int = 12_000) -> tuple:
        """
        Genera (system_msg, user_msg) para evaluación batch de requisitos.
        Formato JSON de respuesta:
        [
          {"presente":"Si","observacion":"cómo se evidencia","confidence":0.95},
          {"presente":"No","observacion":"qué falta","sugerencia":"acción concreta","confidence":0.80}
        ]
        confidence: 0.0–1.0 — certeza del modelo sobre su evaluación.
        """
        reqs_list = "\n".join(f"{i+1}. {p['sub_seccion']}" for i, p in enumerate(params))
        system_msg = (
            "Eres un evaluador académico universitario experto en revisión de material pedagógico. "
            "Determina si cada requisito listado está cubierto en el documento proporcionado. "
            "Evalúa con criterio: cuenta como PRESENTE si el tema aparece de forma explícita "
            "(nombrado directamente) o implícita (ejemplos, demostraciones, desarrollo práctico "
            "que evidencia el concepto aunque use terminología distinta). "
            "SOLO responde con un array JSON válido de exactamente "
            f"{len(params)} elementos, sin texto introductorio ni explicación."
        )
        user_msg = (
            f"CONTEXTO: Sección del curso: \"{section_name}\" | Documento: \"{group_name}\"\n\n"
            f"REQUISITOS A EVALUAR ({len(params)}):\n{reqs_list}\n\n"
            f"CONTENIDO DEL DOCUMENTO:\n{doc_snippet[:max_chars]}\n\n"
            "INSTRUCCIONES PARA CADA REQUISITO:\n"
            "- \"presente\": \"Si\" si hay evidencia clara (explícita o implícita). "
            "  \"No\" si el tema está completamente ausente o es tan superficial que no aporta valor.\n"
            "- \"observacion\": Si PRESENTE → cita brevemente dónde o cómo aparece "
            "  (ej: 'Se desarrolla en la sección de introducción mediante el ejemplo de...'). "
            "  Si AUSENTE → describe específicamente qué falta y por qué es importante.\n"
            "- \"sugerencia\": SOLO si AUSENTE → acción concreta y específica en 1 oración "
            "  para que el autor incluya este contenido (ej: 'Agregar una sección que explique X "
            "  con al menos un ejemplo aplicado').\n"
            "- \"confidence\": 0.0–1.0 (certeza de tu evaluación; usa >0.8 solo cuando la "
            "  evidencia es clara y directa).\n\n"
            "FORMATO EXACTO DE RESPUESTA (array JSON, sin texto extra):\n"
            '[{"presente":"Si","observacion":"Se desarrolla en la introducción con el ejemplo de listas enlazadas","confidence":0.92},\n'
            ' {"presente":"No","observacion":"No se aborda la complejidad temporal de los algoritmos","sugerencia":"Incluir una tabla comparativa de complejidad O(n) para cada algoritmo presentado","confidence":0.88}]\n\n'
            f"Responde SOLO el array JSON con {len(params)} elementos:"
        )
        return system_msg, user_msg

    @staticmethod
    def _parse_batch_item(item: dict, fallback_observacion: str = '') -> dict:
        """Normaliza un elemento de respuesta batch al formato interno."""
        presente = str(item.get('presente', 'No'))
        raw_conf = item.get('confidence', None)
        try:
            confidence = round(float(raw_conf), 3) if raw_conf is not None else None
        except (TypeError, ValueError):
            confidence = None
        return {
            'presente':    presente,
            'observacion': str(item.get('observacion', fallback_observacion)),
            'sugerencia':  str(item.get('sugerencia', '')) if presente != 'Si' else '',
            'confidence':  confidence,
        }

    def _reload_keys_from_db(self, db, user_id=None) -> None:
        """
        Carga las API keys para esta validación con la siguiente prioridad:
          1. Keys personales del usuario (si user_id dado y tiene keys)
          2. Keys del sistema en BD  (settings admin)
          3. Keys del entorno (.env)
        Las keys personales SUSTITUYEN completamente a las del sistema/entorno
        para ese proveedor — garantizando que cada usuario usa solo sus propias keys.
        """
        if db is None:
            return
        # Limpiar exhausted_keys al inicio de cada validación para que las keys
        # cargadas desde BD siempre se intenten — el estado de 429 de una sesión
        # anterior no debe bloquear la siguiente validación.
        self._exhausted_keys.clear()
        try:
            from app.services.settings_service import settings_service
            from app.models.user import User

            # Leer keys personales del usuario
            user_personal = {}
            if user_id is not None:
                try:
                    u = db.query(User).filter(User.id == user_id).first()
                    user_personal = (getattr(u, "personal_api_keys", None) or {}) if u else {}
                except Exception:
                    pass

            def _prefer_personal(provider, system_keys, env_key):
                personal = [k.strip() for k in user_personal.get(provider, []) if k.strip()]
                if personal:
                    print(f"    🔑 Usando {len(personal)} key(s) personal(es) de usuario para {provider}")
                    return personal, "personal"
                # Admin (BD) primero, .env como último recurso
                merged = []
                for k in system_keys:
                    k = k.strip()
                    if k and k not in merged:
                        merged.append(k)
                # env_key puede ser str o list
                env_list = env_key if isinstance(env_key, list) else ([env_key] if env_key else [])
                for k in env_list:
                    if k and k not in merged:
                        merged.append(k)
                admin_keys = [k.strip() for k in system_keys if k.strip()]
                if not merged:
                    return [], "none"
                return merged, "admin" if admin_keys else "env"

            # ── Guardar env keys originales para merge (no usar self.* que ya fue mutado)
            env_gemini    = list(self._default_api_keys)   # todas las keys .env de Gemini
            env_deepseek  = getattr(settings, 'DEEPSEEK_API_KEY', '') or ''
            env_groq      = getattr(settings, 'GROQ_API_KEY', '') or ''
            env_openrouter = getattr(settings, 'OPENROUTER_API_KEY', '') or ''

            # ── Gemini ──────────────────────────────────────────────────────────
            db_gemini = settings_service.get_json("gemini_api_keys", db, default=[]) or []
            merged_gemini, gemini_src = _prefer_personal("gemini", db_gemini, env_gemini)
            self._api_keys = merged_gemini
            self.enabled = len(merged_gemini) > 0
            if merged_gemini:
                import google.generativeai as genai
                genai.configure(api_key=merged_gemini[0])

            # ── DeepSeek ────────────────────────────────────────────────────────
            db_deepseek = settings_service.get_json("deepseek_api_keys", db, default=[]) or []
            merged_deepseek, deepseek_src = _prefer_personal("deepseek", db_deepseek, env_deepseek or None)
            self._deepseek_key = merged_deepseek[0] if merged_deepseek else None

            # ── Groq ────────────────────────────────────────────────────────────
            db_groq = settings_service.get_json("groq_api_keys", db, default=[]) or []
            merged_groq, groq_src = _prefer_personal("groq", db_groq, env_groq or None)
            self._groq_key = merged_groq[0] if merged_groq else None

            # ── OpenRouter ──────────────────────────────────────────────────────
            db_openrouter = settings_service.get_json("openrouter_api_keys", db, default=[]) or []
            merged_openrouter, openrouter_src = _prefer_personal("openrouter", db_openrouter, env_openrouter or None)
            self._openrouter_key = merged_openrouter[0] if merged_openrouter else None

            # ── Fuente del proveedor primario (para AnalysisLog) ─────────────────
            if merged_gemini:
                self._key_source = gemini_src
            elif self._deepseek_key:
                self._key_source = deepseek_src
            elif self._groq_key:
                self._key_source = groq_src
            elif self._openrouter_key:
                self._key_source = openrouter_src
            else:
                self._key_source = "none"

            print(
                f"  🔑 Keys cargadas — Gemini:{len(merged_gemini)} "
                f"DeepSeek:{'✓' if self._deepseek_key else '✗'} "
                f"Groq:{'✓' if self._groq_key else '✗'} "
                f"OpenRouter:{'✓' if self._openrouter_key else '✗'}"
            )

        except Exception as e:
            import traceback
            print(f"  ⚠️  Error recargando keys desde BD: {e}")
            print(traceback.format_exc())

    @property
    def primary_provider(self) -> str:
        if self.enabled:
            return "gemini"
        if self._deepseek_key:
            return "deepseek"
        if self._groq_key:
            return "groq"
        if self._openrouter_key:
            return "openrouter"
        return "basic"

    @property
    def key_source(self) -> str:
        return self._key_source

    def _get_available_keys(self) -> List[str]:
        """Retorna las claves Gemini que aún no han sido marcadas como agotadas."""
        return [k for k in self._api_keys if k not in self._exhausted_keys]

    def _call_gemini_raw(self, api_key: str, model_name: str, prompt: str) -> str:
        """
        Llama a Gemini con una clave específica y retorna el texto crudo.
        Configura la clave antes de llamar (cada llamada puede usar una clave diferente).
        Lanza excepciones — el llamador decide cómo manejarlas (rotar, fallback, etc.).
        """
        import google.generativeai as genai
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(model_name)
        response = model.generate_content(prompt)
        raw = response.text.strip()
        raw = re.sub(r'^```json?\s*', '', raw)
        raw = re.sub(r'\s*```$', '', raw)
        return raw

    def _call_groq_batch(
        self,
        params: List[Dict],
        doc_snippet: str,
        section_name: str,
        group_name: str,
    ) -> Optional[List[Dict[str, str]]]:
        """
        Llama a Groq (proveedor alternativo gratuito, compatible con OpenAI) para
        evaluar todos los parámetros del grupo.  No requiere dependencias extras —
        usa urllib.request de la biblioteca estándar de Python.

        Modelos Groq actuales (2026):
          - llama-3.1-8b-instant  (rápido, ligero — principal)
          - llama-3.3-70b-versatile (más capaz — fallback)
        """
        import urllib.request
        import urllib.error

        system_msg, user_msg = self._batch_prompt(params, section_name, group_name, doc_snippet, 12_000)

        # Intentar con modelos Groq en orden de preferencia
        groq_models = ["llama-3.1-8b-instant", "llama-3.3-70b-versatile"]

        for groq_model in groq_models:
            payload = json.dumps({
                "model": groq_model,
                "messages": [
                    {"role": "system", "content": system_msg},
                    {"role": "user",   "content": user_msg},
                ],
                "temperature": getattr(self, '_req_temperature', 0.05),
                "max_tokens":  getattr(self, '_req_max_tokens', 2000),
            }).encode("utf-8")

            try:
                req = urllib.request.Request(
                    "https://api.groq.com/openai/v1/chat/completions",
                    data=payload,
                    headers={
                        "Authorization": f"Bearer {self._groq_key}",
                        "Content-Type": "application/json",
                        "User-Agent": "python-httpx/0.27.0",
                    },
                    method="POST",
                )
                with urllib.request.urlopen(req, timeout=60) as resp:
                    result = json.loads(resp.read().decode("utf-8"))
                    content = result["choices"][0]["message"]["content"].strip()
                    content = re.sub(r"^```json?\s*", "", content)
                    content = re.sub(r"\s*```$", "", content)

                    parsed = json.loads(content)
                    if isinstance(parsed, dict):
                        for v in parsed.values():
                            if isinstance(v, list):
                                parsed = v
                                break

                    if not isinstance(parsed, list):
                        print(f"  ⚠️  Groq ({groq_model}): respuesta no es lista")
                        continue

                    results = []
                    for i, p in enumerate(params):
                        if i < len(parsed):
                            parsed_item = self._parse_batch_item(parsed[i], '')
                            parsed_item['observacion'] = f"[Groq] {parsed_item['observacion']}"
                            results.append(parsed_item)
                        else:
                            results.append(self._fallback_keyword_check(p["sub_seccion"], doc_snippet))

                    print(f"  🤖 [MODELO: GROQ {groq_model}] respondió — {len(results)} resultados")
                    return results

            except urllib.error.HTTPError as e:
                body = ""
                try:
                    body = e.read().decode("utf-8")[:200]
                except Exception:
                    pass
                if e.code == 429:
                    print(f"  ⚠️  Groq ({groq_model}) quota agotada")
                elif "decommissioned" in body or "deprecated" in body:
                    print(f"  ⚠️  Groq ({groq_model}) descontinuado, probando siguiente")
                    continue
                else:
                    print(f"  ⚠️  Groq ({groq_model}) HTTP {e.code}: {body[:80]}")
            except json.JSONDecodeError:
                print(f"  ⚠️  Groq ({groq_model}): JSON inválido, probando siguiente")
                continue
            except Exception as e:
                print(f"  ⚠️  Groq ({groq_model}): {e}")

        return None

    def _call_deepseek_batch(
        self,
        params: List[Dict],
        doc_snippet: str,
        section_name: str,
        group_name: str,
    ) -> Optional[List[Dict[str, str]]]:
        """
        Llama a DeepSeek (compatible con OpenAI) para evaluar todos los
        parámetros del grupo en una sola llamada.
        Modelo: deepseek-chat (DeepSeek-V3) — rápido, barato, preciso.
        """
        import urllib.request
        import urllib.error

        system_msg, user_msg = self._batch_prompt(params, section_name, group_name, doc_snippet, 15_000)

        payload = json.dumps({
            "model": "deepseek-chat",
            "messages": [
                {"role": "system", "content": system_msg},
                {"role": "user",   "content": user_msg},
            ],
            "temperature": getattr(self, '_req_temperature', 0.05),
            "max_tokens":  getattr(self, '_req_max_tokens', 2000),
        }).encode("utf-8")

        try:
            req = urllib.request.Request(
                "https://api.deepseek.com/v1/chat/completions",
                data=payload,
                headers={
                    "Authorization": f"Bearer {self._deepseek_key}",
                    "Content-Type": "application/json",
                },
                method="POST",
            )
            with urllib.request.urlopen(req, timeout=60) as resp:
                data = json.loads(resp.read().decode("utf-8"))

            raw = data["choices"][0]["message"]["content"].strip()
            raw = re.sub(r'^```json?\s*', '', raw)
            raw = re.sub(r'\s*```$', '', raw)
            parsed = json.loads(raw)

            if not isinstance(parsed, list):
                print("  ⚠️  DeepSeek: respuesta no es lista")
                return None

            results = []
            for i, p in enumerate(params):
                if i < len(parsed):
                    results.append(self._parse_batch_item(parsed[i]))
                else:
                    results.append(self._fallback_keyword_check(p["sub_seccion"], doc_snippet))
            print(f"  🤖 [MODELO: DeepSeek-V3] respondió — {len(results)} resultados")
            return results

        except urllib.error.HTTPError as e:
            body = e.read().decode("utf-8", errors="replace")[:120]
            print(f"  ⚠️  DeepSeek HTTP {e.code}: {body}")
        except json.JSONDecodeError:
            print("  ⚠️  DeepSeek: JSON inválido")
        except Exception as e:
            print(f"  ⚠️  DeepSeek: {e}")

        return None

    def _call_openrouter_batch(
        self,
        params: List[Dict],
        doc_snippet: str,
        section_name: str,
        group_name: str,
    ) -> Optional[List[Dict[str, str]]]:
        """
        Llama a OpenRouter (proveedor alternativo gratuito, compatible con OpenAI).
        Modelos gratuitos disponibles (etiqueta :free en openrouter.ai):
          - meta-llama/llama-3.1-8b-instruct:free
          - google/gemma-3-12b-it:free
        """
        import urllib.request
        import urllib.error

        system_msg, user_msg = self._batch_prompt(params, section_name, group_name, doc_snippet, 12_000)

        openrouter_models = [
            "meta-llama/llama-3.1-8b-instruct:free",
            "google/gemma-3-12b-it:free",
        ]

        for model in openrouter_models:
            payload = json.dumps({
                "model": model,
                "messages": [
                    {"role": "system", "content": system_msg},
                    {"role": "user",   "content": user_msg},
                ],
                "temperature": getattr(self, '_req_temperature', 0.05),
                "max_tokens":  getattr(self, '_req_max_tokens', 2000),
            }).encode("utf-8")

            try:
                req = urllib.request.Request(
                    "https://openrouter.ai/api/v1/chat/completions",
                    data=payload,
                    headers={
                        "Authorization": f"Bearer {self._openrouter_key}",
                        "Content-Type": "application/json",
                        "User-Agent": "python-httpx/0.27.0",
                        "HTTP-Referer": "https://cococys.app",
                    },
                    method="POST",
                )
                with urllib.request.urlopen(req, timeout=60) as resp:
                    result = json.loads(resp.read().decode("utf-8"))
                    content = result["choices"][0]["message"]["content"].strip()
                    content = re.sub(r"^```json?\s*", "", content)
                    content = re.sub(r"\s*```$", "", content)

                    parsed = json.loads(content)
                    if isinstance(parsed, dict):
                        for v in parsed.values():
                            if isinstance(v, list):
                                parsed = v
                                break

                    if not isinstance(parsed, list):
                        print(f"  ⚠️  OpenRouter ({model}): respuesta no es lista")
                        continue

                    results = []
                    for i, p in enumerate(params):
                        if i < len(parsed):
                            parsed_item = self._parse_batch_item(parsed[i])
                            parsed_item['observacion'] = f"[OpenRouter] {parsed_item['observacion']}"
                            results.append(parsed_item)
                        else:
                            results.append(self._fallback_keyword_check(p["sub_seccion"], doc_snippet))

                    print(f"  🤖 [MODELO: OPENROUTER {model}] respondió — {len(results)} resultados")
                    return results

            except urllib.error.HTTPError as e:
                body = ""
                try:
                    body = e.read().decode("utf-8")[:200]
                except Exception:
                    pass
                if e.code == 429:
                    print(f"  ⚠️  OpenRouter ({model}) quota agotada")
                elif e.code == 402:
                    print(f"  ⚠️  OpenRouter ({model}) requiere créditos — probando siguiente")
                    continue
                else:
                    print(f"  ⚠️  OpenRouter ({model}) HTTP {e.code}: {body[:80]}")
            except json.JSONDecodeError:
                print(f"  ⚠️  OpenRouter ({model}): JSON inválido, probando siguiente")
                continue
            except Exception as e:
                print(f"  ⚠️  OpenRouter ({model}): {e}")

        return None

    # ──────────────────────────────────────────────────────────────────────────
    # Normalización (igual que structure_validation_service)
    # ──────────────────────────────────────────────────────────────────────────

    def _normalize(self, name: str) -> str:
        name = re.sub(r'^\d+[\s_\-\.]+', '', str(name).strip())
        name = name.replace('_', ' ').replace('-', ' ')
        name = re.sub(r'\s+', ' ', name).strip()
        name = name.lower()
        nfkd = unicodedata.normalize('NFKD', name)
        return ''.join(c for c in nfkd if not unicodedata.combining(c))

    def _derive_section_from_folder(self, folder_name: str) -> str:
        """
        Convierte un nombre de carpeta Drive en el nombre canónico de sección.
        "Semana_2" → "Semana 2", "semana_12" → "Semana 12"
        Para otros tipos devuelve el nombre limpio (underscores → espacios).
        """
        normalized = self._normalize(folder_name)
        m = re.match(r'^semana\s*(\d+)', normalized)
        if m:
            return f"Semana {m.group(1)}"
        return folder_name.replace('_', ' ').replace('-', ' ').strip()

    def _get_folder_type(self, folder_name: str) -> str:
        """
        Detecta el tipo de carpeta.
        Retorna: 'semana' | 'proyectos' | 'practicas' | 'tareas' | 'unknown'
        """
        norm = self._normalize(folder_name)
        if re.match(r'^semana\s*\d+', norm):
            return 'semana'
        for key in LAB_FOLDER_TYPES:
            if key in norm:
                return key
        return 'unknown'

    def get_lab_sheet(self, wb: openpyxl.Workbook, folder_type: str):
        """
        Obtiene la hoja del Excel correspondiente a una carpeta de laboratorio.
        folder_type: 'proyectos' | 'practicas' | 'tareas'
        Busca por nombre exacto, luego por coincidencia parcial.
        """
        sheet_name = LAB_FOLDER_TYPES.get(folder_type, '')
        # Búsqueda exacta (case-insensitive)
        for name in wb.sheetnames:
            if self._normalize(name) == self._normalize(sheet_name):
                print(f"  📋 Hoja lab '{name}' encontrada")
                return wb[name]
        # Búsqueda parcial
        for name in wb.sheetnames:
            if sheet_name.lower() in name.lower():
                print(f"  📋 Hoja lab '{name}' encontrada (parcial)")
                return wb[name]
        print(f"  ⚠️  No se encontró hoja '{sheet_name}'. Hojas disponibles: {wb.sheetnames}")
        return None

    # ──────────────────────────────────────────────────────────────────────────
    # Excel: localizar, leer y parsear
    # ──────────────────────────────────────────────────────────────────────────

    REVISION_FOLDER = "revision de material"

    def _find_revision_subfolder(self, folder_id: str) -> Optional[str]:
        """Busca la carpeta '0. Revision de Material' (o similar) dentro de folder_id."""
        try:
            subfolders = drive_service.list_folders(folder_id)
            for sf in subfolders:
                if self.REVISION_FOLDER in self._normalize(sf.get('name', '')):
                    return sf['id']
        except Exception:
            pass
        return None

    def find_matrix_file(self, folder_id: str) -> Optional[Tuple[str, str]]:
        """
        Busca el Excel 'Matriz observaciones*.xlsx' en la carpeta indicada.
        Primero busca dentro de '0. Revision de Material'; si no, busca en folder_id.
        Retorna (file_id, file_name) o None.
        """
        try:
            search_id = self._find_revision_subfolder(folder_id) or folder_id
            files = drive_service.list_files(search_id)
            names = [f.get('name', '') for f in files]
            print(f"  📂 Archivos en {search_id}: {names}")
            for f in files:
                name = f.get('name', '')
                if name.lower().startswith(MATRIX_FILE_PREFIX.lower()):
                    return (f['id'], name)
            # Fallback: si buscamos en subfolder y no encontramos, intentar en raíz
            if search_id != folder_id:
                files = drive_service.list_files(folder_id)
                for f in files:
                    name = f.get('name', '')
                    if name.lower().startswith(MATRIX_FILE_PREFIX.lower()):
                        return (f['id'], name)
            return None
        except Exception as e:
            print(f"Error buscando matriz: {e}")
            return None

    def read_workbook(self, file_id: str) -> Optional[openpyxl.Workbook]:
        """Descarga el Excel desde Drive y lo carga con openpyxl."""
        content = drive_service.download_file(file_id)
        if not content:
            print("No se pudo descargar el Excel")
            return None
        try:
            return openpyxl.load_workbook(io.BytesIO(content))
        except Exception as e:
            print(f"Error cargando workbook: {e}")
            return None

    def get_content_sheet(self, wb: openpyxl.Workbook) -> Optional[Any]:
        """
        Obtiene la hoja 'Matriz observaciones' (2da hoja por índice).
        Fallback: busca por nombre que contenga 'observaciones'.
        """
        if len(wb.sheetnames) >= 2:
            return wb.worksheets[1]
        for name in wb.sheetnames:
            if 'observacion' in self._normalize(name):
                return wb[name]
        print(f"⚠️  No se encontró la hoja 'Matriz observaciones'. Hojas disponibles: {wb.sheetnames}")
        return None

    def _find_header_row(self, ws, min_matches: int = 3) -> Tuple[int, Dict[str, int]]:
        """
        Escanea las primeras 5 filas para encontrar la fila de encabezados.
        Retorna (header_row_idx, col_map) con índices 1-based.
        col_map keys: 'seccion', 'sub_seccion', 'aplica', 'autor', 'presente', 'observaciones'

        min_matches=2 permite detectar hojas lab con menos columnas estándar.
        """
        target_keywords = {
            'seccion', 'sub seccion', 'aplica', 'autor', 'presente', 'observaciones'
        }

        best_row: Optional[Tuple[int, Dict, int]] = None  # (row_idx, col_map, match_count)

        for row_idx in range(1, 6):
            row_vals = []
            for col_idx in range(1, ws.max_column + 1):
                val = ws.cell(row=row_idx, column=col_idx).value
                row_vals.append(self._normalize(str(val or '')))

            matches = sum(1 for v in row_vals if any(kw in v for kw in target_keywords))
            if matches >= min_matches:
                col_map: Dict[str, int] = {}
                for col_idx, norm_val in enumerate(row_vals, 1):
                    if not norm_val:
                        continue
                    if 'sub' in norm_val and 'seccion' in norm_val:
                        col_map['sub_seccion'] = col_idx
                    elif 'seccion' in norm_val:
                        col_map['seccion'] = col_idx
                    elif norm_val.startswith('aplica'):
                        col_map['aplica'] = col_idx
                    elif norm_val == 'autor':
                        col_map['autor'] = col_idx
                    elif 'presente' in norm_val:
                        col_map['presente'] = col_idx
                    elif 'observaci' in norm_val:
                        col_map['observaciones'] = col_idx
                print(f"📋 Encabezados encontrados en fila {row_idx}: {col_map}")
                return (row_idx, col_map)
            # Guardar la mejor fila aunque no alcance el mínimo
            if best_row is None or matches > best_row[2]:
                col_map_partial: Dict[str, int] = {}
                for col_idx, norm_val in enumerate(row_vals, 1):
                    if not norm_val:
                        continue
                    if 'sub' in norm_val and 'seccion' in norm_val:
                        col_map_partial['sub_seccion'] = col_idx
                    elif 'seccion' in norm_val:
                        col_map_partial['seccion'] = col_idx
                    elif norm_val.startswith('aplica'):
                        col_map_partial['aplica'] = col_idx
                    elif norm_val == 'autor':
                        col_map_partial['autor'] = col_idx
                    elif 'presente' in norm_val:
                        col_map_partial['presente'] = col_idx
                    elif 'observaci' in norm_val:
                        col_map_partial['observaciones'] = col_idx
                best_row = (row_idx, col_map_partial, matches)

        raise ValueError("No se encontró fila de encabezados en las primeras 5 filas del Excel")

    def parse_requirements(
        self,
        ws,
        target_section: str
    ) -> Tuple[List[Dict], Dict[str, int], int]:
        """
        Lee la hoja y filtra los requisitos para la sección indicada donde Aplica == 'Si'.
        Retorna (requirements_list, col_map, header_row_idx).

        Cada elemento de requirements_list:
            {'row_idx': int, 'sub_seccion': str, 'autor': str}
        """
        header_row_idx, col_map = self._find_header_row(ws)
        norm_target = self._normalize(target_section)
        requirements = []
        last_sec_val = None   # soporte para celdas combinadas (merged cells)

        for row_idx in range(header_row_idx + 1, ws.max_row + 1):
            # Leer sección
            sec_col = col_map.get('seccion')
            if not sec_col:
                continue
            sec_val = ws.cell(row=row_idx, column=sec_col).value
            if sec_val:
                last_sec_val = sec_val
            elif last_sec_val:
                sec_val = last_sec_val
            if not sec_val:
                continue
            if self._normalize(str(sec_val)) != norm_target:
                continue

            # Verificar Aplica == Si
            aplica_col = col_map.get('aplica')
            if aplica_col:
                aplica_val = ws.cell(row=row_idx, column=aplica_col).value
                if self._normalize(str(aplica_val or '')) != 'si':
                    continue

            # Leer sub-sección
            sub_col = col_map.get('sub_seccion')
            sub_val = ws.cell(row=row_idx, column=sub_col).value if sub_col else None
            if not sub_val or str(sub_val).strip() == '':
                continue

            # Leer autor (opcional)
            autor_col = col_map.get('autor')
            autor_val = ws.cell(row=row_idx, column=autor_col).value if autor_col else None

            requirements.append({
                'row_idx':    row_idx,
                'sub_seccion': str(sub_val).strip(),
                'autor':      str(autor_val).strip() if autor_val else '',
            })

        print(f"📋 Requisitos para '{target_section}': {len(requirements)}")

        # DEBUG: mostrar primeras 10 filas de datos para diagnóstico
        if len(requirements) == 0:
            sec_col = col_map.get('seccion')
            aplica_col = col_map.get('aplica')
            print(f"🔍 DEBUG: buscando sección normalizada = '{norm_target}'")
            seen_sections = set()
            for r in range(header_row_idx + 1, min(header_row_idx + 30, ws.max_row + 1)):
                sec_raw = ws.cell(row=r, column=sec_col).value if sec_col else None
                aplica_raw = ws.cell(row=r, column=aplica_col).value if aplica_col else None
                sec_norm = self._normalize(str(sec_raw or ''))
                if sec_norm and sec_norm not in seen_sections:
                    seen_sections.add(sec_norm)
                    print(f"   Fila {r}: sección='{sec_raw}' (norm='{sec_norm}') | aplica='{aplica_raw}'")
            if not seen_sections:
                print("   ⚠️  No se encontraron valores en la columna Sección")

        return requirements, col_map, header_row_idx

    # ──────────────────────────────────────────────────────────────────────────
    # Extracción de texto
    # ──────────────────────────────────────────────────────────────────────────

    def _extract_pdf_text(self, file_bytes: bytes) -> str:
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            pages = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ''
                if text.strip():
                    pages.append(f"=== Página {i} ===\n{text}")
            return "\n".join(pages)
        except Exception as e:
            print(f"Error extrayendo PDF: {e}")
            return ''

    def _extract_docx_text(self, file_bytes: bytes) -> str:
        try:
            doc = DocxDocument(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs)
        except Exception as e:
            print(f"Error extrayendo DOCX: {e}")
            return ''

    def _extract_pptx_text(self, file_bytes: bytes) -> str:
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                parts = [f"=== Diapositiva {i} ==="]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                slides_text.append("\n".join(parts))
            return "\n\n".join(slides_text)
        except Exception as e:
            print(f"Error extrayendo PPTX: {e}")
            return ''

    def extract_text_from_bytes(self, file_bytes: bytes, mime_type: str) -> str:
        """Extrae texto según el tipo MIME del archivo."""
        if mime_type == 'application/pdf':
            return self._extract_pdf_text(file_bytes)
        elif 'wordprocessingml' in mime_type:
            return self._extract_docx_text(file_bytes)
        elif 'presentationml' in mime_type:
            return self._extract_pptx_text(file_bytes)
        return ''

    def parse_requirements_by_group(
        self,
        ws,
        target_section: Optional[str],
        col_map: Dict[str, int],
        header_row_idx: int
    ) -> List[Dict]:
        """
        Lee la hoja y agrupa los requisitos por tipo de documento (Presentación, Lectura, Video…).

        target_section=None → lee TODOS los requisitos sin filtrar por sección.
        Útil para hojas de Proyectos/Practicas/Tareas donde toda la hoja pertenece
        al mismo tipo de carpeta.

        Estructura del Excel:
          Sección  | Sub-sección          | Aplica
          Semana 2 | Presentación         |        ← encabezado de grupo (sin Aplica)
          Semana 2 | Bienvenida           | Si     ← parámetro
          Semana 2 | Lectura              |        ← encabezado de grupo
          Semana 2 | Titulo               | Si     ← parámetro
        """
        norm_target = self._normalize(target_section) if target_section else None
        sec_col    = col_map.get('seccion')
        sub_col    = col_map.get('sub_seccion')
        aplica_col = col_map.get('aplica')
        autor_col  = col_map.get('autor')

        # Para hojas lab de una sola columna (Proyectos/Practicas/Tareas), la columna
        # "Sección" contiene tanto los encabezados de grupo ("Proyecto 1") como los
        # criterios ("Titulo Proyecto"). Si no hay sub_seccion, usar seccion como fallback.
        effective_sub_col = sub_col or sec_col
        if not effective_sub_col:
            return []

        groups: List[Dict] = []
        current_group: Optional[Dict] = None
        last_sec_val = None

        for row_idx in range(header_row_idx + 1, ws.max_row + 1):
            # Filtro por sección (si aplica)
            if norm_target is not None and sec_col:
                sec_val = ws.cell(row=row_idx, column=sec_col).value
                if sec_val:
                    last_sec_val = sec_val
                elif last_sec_val:
                    sec_val = last_sec_val
                if not sec_val:
                    continue
                if self._normalize(str(sec_val)) != norm_target:
                    continue

            sub_val = ws.cell(row=row_idx, column=effective_sub_col).value
            if not sub_val or str(sub_val).strip() == '':
                continue
            sub_str = str(sub_val).strip()

            aplica_val = ws.cell(row=row_idx, column=aplica_col).value if aplica_col else None
            aplica_norm = self._normalize(str(aplica_val or ''))

            # Excluir filas explícitamente marcadas como No aplica
            if aplica_norm == 'no':
                continue

            sub_norm = self._normalize(sub_str)

            # ── Detección de encabezado de grupo ────────────────────────────
            # Criterio 1: vocabulario de tipos de documento (Semana sheets)
            # Criterio 2: patrón lab "Proyecto N", "Practica N", "Tarea N"
            _is_lab_header = bool(
                re.match(r'^(proyecto|practica|tarea)\s*\d*$', sub_norm)
            )

            if sub_norm in DOCUMENT_TYPE_HEADERS or _is_lab_header:
                # Tipo de documento o encabezado de lab → nuevo grupo
                current_group = {'group_name': sub_str, 'params': []}
                groups.append(current_group)
            else:
                # Parámetro a evaluar
                autor_val = ws.cell(row=row_idx, column=autor_col).value if autor_col else None
                param = {
                    'row_idx':    row_idx,
                    'sub_seccion': sub_str,
                    'autor':      str(autor_val).strip() if autor_val else '',
                }
                if current_group is not None:
                    current_group['params'].append(param)
                else:
                    # No hay encabezado previo → grupo implícito "General"
                    current_group = {'group_name': 'General', 'params': [param]}
                    groups.append(current_group)

        # Descartar grupos vacíos
        groups = [g for g in groups if g['params']]

        total_params = sum(len(g['params']) for g in groups)
        print(f"📋 Grupos para '{target_section}': {len(groups)} grupos, {total_params} parámetros")
        for g in groups:
            print(f"   📑 '{g['group_name']}': {len(g['params'])} parámetros")
            for p in g['params']:
                print(f"      · {p['sub_seccion']}")

        return groups

    def _parse_lab_sheet_raw(self, ws) -> List[Dict]:
        """
        Fallback robusto para hojas lab con estructura no estándar.
        Lee la primera columna significativa de cada fila.
        Detecta encabezados de grupo ("Proyecto 1", "Practica 2"…) y crea
        grupos separados — igual que parse_requirements_by_group pero sin
        depender del mapeo de columnas.
        """
        SKIP_NORMS = {
            'si', 'no', 'aplica', 'sub seccion', 'seccion', 'autor',
            'presente', 'observaciones', 'observacion', 'criterio',
            'criterios', 'requisito', 'requisitos', 'nombre', 'fecha',
            'calificacion', 'puntaje', 'nota', 'semana', 'sub-seccion',
            'descripcion', 'descripcion',
        }

        groups: List[Dict] = []
        current_group: Optional[Dict] = None
        seen_params: set = set()

        for row_idx in range(1, min(ws.max_row + 1, 500)):
            # Leer la primera celda no vacía de la fila (columna más a la izquierda)
            row_text = None
            row_col  = None
            for col_idx in range(1, min(ws.max_column + 1, 5)):
                val = ws.cell(row=row_idx, column=col_idx).value
                if val and str(val).strip():
                    row_text = str(val).strip()
                    row_col  = col_idx
                    break
            if not row_text:
                continue

            norm = self._normalize(row_text)

            # Saltar palabras clave de encabezado de tabla
            if norm in SKIP_NORMS or len(row_text) < 4:
                continue

            # ── Encabezado de grupo lab ("Proyecto 1", "Practica 2"…) ──
            if re.match(r'^(proyecto|practica|tarea)\s*\d*$', norm):
                current_group = {'group_name': row_text, 'params': []}
                groups.append(current_group)
                continue

            # ── Criterio de evaluación ──────────────────────────────────
            if len(row_text) < 8 or norm in seen_params:
                continue
            seen_params.add(norm)

            if current_group is None:
                current_group = {'group_name': 'General', 'params': []}
                groups.append(current_group)

            current_group['params'].append({
                'row_idx':    row_idx,
                'sub_seccion': row_text,
                'autor':      '',
            })

        groups = [g for g in groups if g['params']]
        if groups:
            total = sum(len(g['params']) for g in groups)
            print(f"  📋 Parsing raw: {len(groups)} grupos, {total} requisitos")
        return groups

    def _parse_lab_sheet_full(self, ws) -> Tuple[List[Dict], Dict[str, int]]:
        """
        Parser dedicado y completo para hojas Proyectos/Practicas/Tareas.

        Hace en un solo paso:
          1. Escanea toda la hoja para localizar las columnas de escritura
             (Presente, Observaciones) — sin depender de _find_header_row.
          2. Detecta la columna de criterios (primera columna con contenido).
          3. Detecta la columna Aplica (si existe).
          4. Agrupa criterios bajo encabezados "Proyecto N / Practica N / Tarea N".
          5. Retorna (groups, write_cols) donde write_cols tiene las claves
             'presente' y/o 'observaciones' con sus índices de columna 1-based.

        Esto desacopla completamente el parsing lab del flujo de hojas Semana.
        """
        # ── Paso 1: localizar columnas de escritura en toda la hoja ──────────
        write_cols: Dict[str, int] = {}
        aplica_col: Optional[int] = None
        autor_col:  Optional[int] = None
        criteria_col: Optional[int] = None   # columna principal con texto de criterios

        # Escaneamos las primeras 15 filas buscando encabezados de columna
        for r in range(1, min(ws.max_row + 1, 15)):
            row_matched = False
            for c in range(1, min(ws.max_column + 1, 30)):
                raw = ws.cell(row=r, column=c).value
                norm = self._normalize(str(raw or ''))
                if not norm:
                    continue
                if 'presente' in norm and 'presente' not in write_cols:
                    write_cols['presente'] = c
                    row_matched = True
                elif 'observaci' in norm and 'observaciones' not in write_cols:
                    write_cols['observaciones'] = c
                    row_matched = True
                elif norm.startswith('aplica') and aplica_col is None:
                    aplica_col = c
                elif norm == 'autor' and autor_col is None:
                    autor_col = c
            if row_matched and len(write_cols) >= 2:
                break   # ambas columnas encontradas

        if autor_col:
            write_cols['autor'] = autor_col
        print(f"  📝 Columnas detectadas: {write_cols} | aplica: {aplica_col} | autor: {autor_col}")

        # ── Paso 2: detectar columna de criterios ────────────────────────────
        # Es la primera columna (col. más a la izquierda) que contiene texto
        # significativo en las primeras 30 filas (excluyendo columnas de escritura).
        write_col_idxs = set(write_cols.values())
        for r in range(1, min(ws.max_row + 1, 30)):
            for c in range(1, min(ws.max_column + 1, 10)):
                if c in write_col_idxs or c == aplica_col:
                    continue
                raw = ws.cell(row=r, column=c).value
                if raw and str(raw).strip() and len(str(raw).strip()) >= 3:
                    norm = self._normalize(str(raw))
                    # Descartar filas de encabezado de tabla
                    if any(kw in norm for kw in ('presente', 'observaci', 'aplica', 'autor')):
                        continue
                    criteria_col = c
                    break
            if criteria_col:
                break

        if criteria_col is None:
            criteria_col = 1   # fallback: columna A
        print(f"  📌 Columna de criterios: {criteria_col}")

        # ── Paso 3: recorrer toda la hoja y agrupar ───────────────────────────
        SKIP_NORMS = {
            'si', 'no', 'aplica', 'sub seccion', 'seccion', 'sub-seccion',
            'presente', 'observaciones', 'observacion', 'autor', 'criterio',
            'criterios', 'requisito', 'requisitos', 'nombre', 'fecha',
            'calificacion', 'puntaje', 'nota', 'descripcion',
        }

        # Columnas que NO son de criterios ni de escritura (candidatas a encabezado de grupo)
        excluded_cols = write_col_idxs | ({aplica_col} if aplica_col else set())
        # Buscar en todas las columnas disponibles, no solo criteria_col
        max_scan_col = min(ws.max_column + 1, 15)

        def _is_group_header(text: str) -> bool:
            """True si el texto corresponde a un encabezado de grupo lab."""
            if not text or len(text.strip()) >= 80:
                return False
            norm = self._normalize(text.strip())
            norm_clean = re.sub(r'[​‌‍﻿­]', '', norm).strip()
            return bool(re.match(
                r'^(proyecto|practica|tarea)\s*\d*(\s*[\-:,./].*)?$',
                norm_clean
            ))

        groups: List[Dict] = []
        current_group: Optional[Dict] = None
        seen_in_group: set = set()   # (group_name, norm_criteria) → evitar duplicados

        for r in range(1, ws.max_row + 1):
            # ── Buscar encabezado de grupo en CUALQUIER columna de la fila ──
            # Esto resuelve el caso donde los encabezados (Proyecto 1, Proyecto 2…)
            # están en una columna distinta a la de los criterios.
            group_header_found = False
            for c in range(1, max_scan_col):
                if c in excluded_cols:
                    continue
                raw_h = ws.cell(row=r, column=c).value
                if raw_h and _is_group_header(str(raw_h)):
                    header_text = str(raw_h).strip()
                    current_group = {'group_name': header_text, 'params': [], 'header_row': r}
                    groups.append(current_group)
                    seen_in_group = set()
                    print(f"     🗂️  Grupo detectado: '{header_text}' (fila {r}, col {c})")
                    group_header_found = True
                    break
            if group_header_found:
                continue

            # ── Leer criterio desde criteria_col ────────────────────────────
            raw = ws.cell(row=r, column=criteria_col).value
            if not raw or not str(raw).strip():
                continue

            text = str(raw).strip()
            norm = self._normalize(text)

            # Saltar encabezados de tabla y valores muy cortos
            if norm in SKIP_NORMS or len(text) < 3:
                continue

            # ── Verificar columna Aplica ──────────────────────────────────────
            if aplica_col:
                aplica_val = ws.cell(row=r, column=aplica_col).value
                if self._normalize(str(aplica_val or '')) == 'no':
                    continue   # explícitamente excluido

            # ── Criterio de evaluación ────────────────────────────────────────
            if len(text) < 5:
                continue

            group_key = current_group['group_name'] if current_group else 'General'
            dedup_key = (group_key, norm)
            if dedup_key in seen_in_group:
                continue
            seen_in_group.add(dedup_key)

            if current_group is None:
                current_group = {'group_name': 'General', 'params': []}
                groups.append(current_group)

            current_group['params'].append({
                'row_idx':     r,
                'sub_seccion': text,
                'autor':       '',
            })

        grupos_vacios = [g for g in groups if not g['params']]
        if grupos_vacios:
            print(f"  ⚠️  {len(grupos_vacios)} grupo(s) sin criterios detectados: "
                  f"{[g['group_name'] for g in grupos_vacios]}")
        groups = [g for g in groups if g['params']]
        total = sum(len(g['params']) for g in groups)
        print(f"  📋 Lab sheet full: {len(groups)} grupos, {total} criterios "
              f"| criteria_col={criteria_col} | write_cols={write_cols}")
        for g in groups:
            rows = [p['row_idx'] for p in g['params']]
            print(f"     📑 '{g['group_name']}': {len(g['params'])} criterios "
                  f"(filas {min(rows)}-{max(rows)})")
            for p in g['params'][:3]:
                print(f"        · [{p['row_idx']}] {p['sub_seccion'][:60]}")
            if len(g['params']) > 3:
                print(f"        … +{len(g['params'])-3} más")

        return groups, write_cols

    # Pistas de tipo: qué MIME buscar para cada nombre de grupo
    _GROUP_TYPE_HINTS: Dict[str, List[str]] = {
        'presentacion':  ['presentationml'],
        'presentación':  ['presentationml'],
        'diapositivas':  ['presentationml'],
        'slides':        ['presentationml'],
        'lectura':       ['pdf', 'wordprocessingml'],
        'reading':       ['pdf', 'wordprocessingml'],
        'documento':     ['pdf', 'wordprocessingml'],
        'video':         ['pdf', 'wordprocessingml'],   # cuestionario/actividad asociada
        'cuestionario':  ['pdf', 'wordprocessingml'],
        'actividad':     ['pdf', 'wordprocessingml'],
        'guia':          ['pdf', 'wordprocessingml'],
        'guía':          ['pdf', 'wordprocessingml'],
        'tarea':         ['pdf', 'wordprocessingml'],
        'informe':       ['pdf', 'wordprocessingml'],
        'reporte':       ['pdf', 'wordprocessingml'],
    }

    def _names_match(self, group_name: str, file_name: str) -> bool:
        """
        Compara el nombre del grupo (tipo de documento) con el nombre del archivo.
        Devuelve True si el nombre normalizado del grupo aparece en el nombre del archivo.

        Ejemplos:
          "Presentación", "Lab pdc1 Presentacion 6 1S2025.pdf"  → True
          "Lectura",       "Lectura_Semana2.pdf"                 → True
          "Video",         "Lab pdc1 Presentacion 6 1S2025.pdf" → False
        """
        norm_group = self._normalize(group_name).lower()
        norm_file  = self._normalize(file_name).lower()
        # Coincidencia si el grupo (o sus primeras palabras) aparece en el nombre del archivo
        return norm_group in norm_file

    def _match_file_to_group(
        self,
        group_name: str,
        files_metadata: List[Dict],
        already_matched: set
    ) -> Optional[Dict]:
        """
        Encuentra el archivo más adecuado para un tipo de documento (grupo).

        Prioridad:
          1. Coincidencia de nombre normalizado (p.ej. "Presentación" ↔ "Presentacion_S2.pptx")
          2. Para grupos lab numerados ("Proyecto 1"): coincidencia por número en nombre de archivo
          3. Pista por MIME type según el nombre del grupo
          4. Cualquier archivo soportado no asignado aún (orden Drive → posición)
        """
        norm_group = self._normalize(group_name)

        # 1. Coincidencia de nombre
        for f in files_metadata:
            if f['name'] in already_matched:
                continue
            if self._names_match(group_name, f['name']):
                return f

        # 2. Para grupos lab numerados ("Proyecto 1", "Practica 2", etc.)
        #    Intentar match por número en el nombre del archivo
        m_num = re.search(r'\d+', group_name)
        is_lab_group = bool(re.match(r'^(proyecto|practica|tarea)\s*\d', norm_group))
        if is_lab_group and m_num:
            num = m_num.group()
            for f in files_metadata:
                if f['name'] in already_matched:
                    continue
                norm_fname = self._normalize(f['name'])
                # Buscar el número como palabra/token en el nombre del archivo
                if re.search(r'(?<![0-9])' + re.escape(num) + r'(?![0-9])', norm_fname):
                    return f

        # 3. Pista por MIME
        for hint_key, mime_frags in self._GROUP_TYPE_HINTS.items():
            if hint_key in norm_group:
                for f in files_metadata:
                    if f['name'] in already_matched:
                        continue
                    mime = f.get('mimeType', '').lower()
                    if any(frag in mime for frag in mime_frags):
                        return f
                break   # intentar solo la primera pista que aplique

        # 4. Fallback posicional ordenado:
        #    Para grupos lab numerados, preferir archivos que contengan la keyword
        #    del tipo (proyecto/practica/tarea) sobre archivos genéricos.
        lab_keyword = None
        for kw in ('proyecto', 'practica', 'tarea'):
            if kw in norm_group:
                lab_keyword = kw
                break

        unmatched = [f for f in files_metadata if f['name'] not in already_matched]
        if lab_keyword:
            # Primero intentar archivos que contengan la keyword del tipo
            keyword_files = [f for f in unmatched
                             if lab_keyword in self._normalize(f['name'])]
            if keyword_files:
                return keyword_files[0]

        if unmatched:
            return unmatched[0]

        return None

    def download_and_extract_all_docs(
        self, folder_id: str
    ) -> Tuple[Dict[str, str], List[Dict]]:
        """
        Descarga y extrae texto de todos los documentos soportados en la carpeta.
        Excluye el Excel de la matriz.

        Retorna:
          doc_texts:      Dict[filename → extracted_text]
          files_metadata: List[{name, mimeType, id}]  (los archivos procesados)
        """
        files = drive_service.list_files(folder_id)
        doc_texts: Dict[str, str] = {}
        files_metadata: List[Dict] = []

        for f in files:
            name = f.get('name', '')
            mime = f.get('mimeType', '')

            # Excluir matriz y reportes generados por el sistema
            if name.startswith(MATRIX_FILE_PREFIX):
                continue
            if name.startswith(self.REPORT_PREFIX):
                continue
            # Solo tipos soportados
            if mime not in SUPPORTED_MIMES:
                print(f"  ↷ Omitiendo '{name}' (MIME: {mime})")
                continue

            # Para Google Workspace nativos, drive_service.download_file() exporta
            # automáticamente (Slides→PPTX, Docs→DOCX). Necesitamos el MIME efectivo
            # (después de exportar) para elegir el extractor correcto.
            effective_mime = drive_service.get_effective_mime(mime)

            tipo = "exportando" if effective_mime != mime else "descargando"
            print(f"  ⬇️  {tipo.capitalize()} '{name}' [{effective_mime.split('.')[-1]}]...")
            file_bytes = drive_service.download_file(f['id'])
            if not file_bytes:
                print(f"  ⚠️  No se pudo obtener '{name}'")
                continue

            text = self.extract_text_from_bytes(file_bytes, effective_mime)
            doc_texts[name] = text
            # Extraer el nombre del propietario del archivo (Drive owners)
            owners = f.get('owners', [])
            owner_name = owners[0].get('displayName', '') if owners else ''
            files_metadata.append({
                'name':      name,
                'mimeType':  mime,
                'id':        f.get('id', ''),
                'owner':     owner_name,
            })
            print(f"  📄 '{name}': {len(text)} chars | propietario: {owner_name or 'N/A'}")

        return doc_texts, files_metadata

    def combine_document_texts(self, doc_texts: Dict[str, str]) -> str:
        """Une todos los textos con separadores claros de documento."""
        parts = []
        for filename, text in doc_texts.items():
            parts.append(f"========== DOCUMENTO: {filename} ==========\n{text}")
        return "\n\n".join(parts)

    # ──────────────────────────────────────────────────────────────────────────
    # Chunking
    # ──────────────────────────────────────────────────────────────────────────

    def chunk_text(self, text: str, chunk_size: int = MAX_CHUNK_CHARS) -> List[str]:
        """
        Divide texto en chunks de máximo chunk_size caracteres.
        Divide por párrafos (\\n\\n) y aplica solapamiento de CHUNK_OVERLAP chars.
        """
        if len(text) <= chunk_size:
            return [text]

        paragraphs = text.split('\n\n')
        chunks: List[str] = []
        current = ''

        for para in paragraphs:
            if len(current) + len(para) + 2 <= chunk_size:
                current = current + '\n\n' + para if current else para
            else:
                if current:
                    chunks.append(current)
                    # Solapamiento: últimos CHUNK_OVERLAP chars del chunk anterior
                    overlap = current[-CHUNK_OVERLAP:] if len(current) > CHUNK_OVERLAP else current
                    current = overlap + '\n\n' + para
                else:
                    # Párrafo individual muy largo: dividir por oraciones
                    sentences = re.split(r'(?<=[.!?])\s+', para)
                    for sent in sentences:
                        if len(current) + len(sent) + 1 <= chunk_size:
                            current = current + ' ' + sent if current else sent
                        else:
                            if current:
                                chunks.append(current)
                            current = sent

        if current:
            chunks.append(current)

        return chunks

    # ──────────────────────────────────────────────────────────────────────────
    # IA: evaluación de requisitos con Gemini
    # ──────────────────────────────────────────────────────────────────────────

    def _call_gemini(
        self,
        requirement_text: str,
        doc_chunk: str,
        section_name: str,
        autor: str,
        model_name: Optional[str] = None
    ) -> Dict[str, str]:
        """
        Llama a Gemini para evaluar si el chunk de documento cubre el requisito.
        Retorna {'presente': 'Si'|'No', 'observacion': str}.
        """
        prompt = f"""Eres un evaluador académico universitario experto. Determina si el contenido del documento cubre el requisito indicado.

SECCIÓN DEL CURSO: {section_name}
REQUISITO A VERIFICAR: {requirement_text}
AUTOR ESPERADO: {autor if autor else 'No especificado'}

CRITERIOS:
- PRESENTE (Si): el documento aborda el tema de forma explícita O implícita (diferente terminología, ejemplos prácticos, demostraciones, o evidencia clara del concepto).
- AUSENTE (No): el tema no aparece, o aparece de forma tan superficial que no puede considerarse cubierto.

REGLAS PARA LA OBSERVACIÓN:
- Si PRESENTE: indica brevemente cómo o dónde se evidencia en el documento.
- Si AUSENTE: indica qué falta específicamente y una recomendación concreta para mejorarlo.

CONTENIDO DEL DOCUMENTO:
{doc_chunk}

Responde ÚNICAMENTE con JSON válido, sin markdown ni texto adicional:
{{"presente": "Si" o "No", "observacion": "Observación específica y accionable de 1-2 oraciones", "confidence": 0.0-1.0}}"""

        target_model = model_name or GEMINI_MODEL

        for key in self._get_available_keys():
            try:
                raw = self._call_gemini_raw(key, target_model, prompt)
                parsed = json.loads(raw)
                raw_conf = parsed.get('confidence', None)
                try:
                    conf = round(float(raw_conf), 3) if raw_conf is not None else None
                except (TypeError, ValueError):
                    conf = None
                return {
                    'presente':    str(parsed.get('presente', 'No')),
                    'observacion': str(parsed.get('observacion', '')),
                    'confidence':  conf,
                }
            except json.JSONDecodeError:
                print(f"  ⚠️  JSON inválido de Gemini")
                break
            except Exception as e:
                err_str = str(e)
                if '429' in err_str or 'quota' in err_str.lower() or 'RESOURCE_EXHAUSTED' in err_str:
                    self._exhausted_keys.add(key)
                    remaining = len(self._get_available_keys())
                    if remaining > 0:
                        print(f"  🔄 Key agotada, rotando ({remaining} restante{'s' if remaining > 1 else ''})")
                        time.sleep(1)
                        continue
                    else:
                        print("  ⚠️  Todas las claves Gemini agotadas")
                        break
                else:
                    print(f"  ⚠️  Error en Gemini: {e}")
                    break

        return self._fallback_keyword_check(requirement_text, doc_chunk)

    def _fallback_keyword_check(self, requirement_text: str, doc_text: str) -> Dict[str, str]:
        """
        Fallback cuando Gemini no está disponible o falla.
        Búsqueda inteligente de palabras clave con matching parcial (stem-like).
        """
        clean_requirement = str(requirement_text).strip().strip('[]')
        norm_req = self._normalize(clean_requirement)
        norm_doc = self._normalize(doc_text)

        # Palabras significativas (≥4 chars), sin stopwords básicas
        STOPWORDS = {'para', 'como', 'esta', 'este', 'esto', 'cion', 'idad', 'mente', 'del', 'los', 'las', 'con', 'que', 'una', 'por', 'sus'}
        keywords = [w for w in norm_req.split() if len(w) >= 4 and w not in STOPWORDS]

        if not keywords:
            return {
                'presente': 'No',
                'observacion': f'Sin IA: no se pudieron evaluar términos en "{clean_requirement}"',
                'confidence': 0.0,
            }

        # Match parcial: cada keyword se busca como substring Y por sus primeros 5 chars (stem)
        found = []
        missing = []
        for kw in keywords:
            stem = kw[:5]
            if kw in norm_doc or stem in norm_doc:
                found.append(kw)
            else:
                missing.append(kw)

        ratio = len(found) / len(keywords)
        presente = 'Si' if ratio >= 0.35 else 'No'

        found_str   = ', '.join(found[:4])   or 'ninguno'
        missing_str = ', '.join(missing[:4]) or 'ninguno'
        obs = (
            f'Sin IA: coincidencias detectadas: {found_str}'
            + (f'; faltantes: {missing_str}' if missing else '')
            + f'. Cobertura {len(found)}/{len(keywords)} términos'
        )
        # Confianza baja — es solo keyword matching, no IA semántica
        return {'presente': presente, 'observacion': obs, 'confidence': round(ratio * 0.6, 3)}

    def _evaluate_requirement(
        self,
        requirement_text: str,
        combined_text: str,
        section_name: str,
        autor: str,
        use_gemini: bool = True,
        model_name: Optional[str] = None
    ) -> Dict[str, str]:
        """
        Evalúa si el texto combinado de todos los documentos cubre el requisito.
        Si el texto es muy largo, evalúa por chunks — Si en cualquier chunk = presente.
        """
        if not use_gemini or not self.enabled:
            return self._fallback_keyword_check(requirement_text, combined_text)

        if len(combined_text) <= MAX_CHUNK_CHARS:
            return self._call_gemini(requirement_text, combined_text, section_name, autor, model_name)

        # Texto largo: evaluar por chunks
        chunks = self.chunk_text(combined_text)
        last_result = {'presente': 'No', 'observacion': ''}
        for i, chunk in enumerate(chunks, 1):
            print(f"    chunk {i}/{len(chunks)}...")
            result = self._call_gemini(requirement_text, chunk, section_name, autor, model_name)
            last_result = result
            if result['presente'] == 'Si':
                return result
            time.sleep(0.3)  # Rate limiting ligero

        return last_result

    def _evaluate_group_batch(
        self,
        params: List[Dict],
        doc_text: str,
        section_name: str,
        group_name: str,
        use_gemini: bool = True,
        model_name: Optional[str] = None,
    ) -> List[Dict[str, str]]:
        """
        Evalúa TODOS los parámetros de un grupo en UNA SOLA llamada a Gemini.
        Reduce N llamadas a 1 por grupo → de 18 llamadas a 3 por Semana.

        Retorna lista ordenada de {'presente', 'observacion'} por parámetro.
        Si Gemini falla, aplica fallback keyword por cada param.
        """
        has_any_provider = (use_gemini or self._deepseek_key or
                            self._groq_key or self._openrouter_key)
        if not has_any_provider:
            print("  ⚠️  Sin proveedores de IA disponibles → keyword fallback")
            return [self._fallback_keyword_check(p['sub_seccion'], doc_text) for p in params]

        # Truncar texto al máximo razonable para un solo prompt
        max_chars = MAX_CHUNK_CHARS
        doc_snippet = doc_text[:max_chars]

        reqs_list = "\n".join(
            f"{i+1}. {p['sub_seccion']}" for i, p in enumerate(params)
        )

        prompt = f"""Eres un evaluador académico universitario experto. Analiza el contenido del documento y determina con precisión si cada requisito está cubierto.

CONTEXTO:
- Sección del curso: {section_name}
- Tipo de documento evaluado: {group_name}

CRITERIOS DE EVALUACIÓN:
- PRESENTE (Si): el documento aborda el tema de forma explícita O implícita (diferente terminología, ejemplos, demostraciones prácticas, o el concepto queda claramente evidenciado).
- AUSENTE (No): el tema NO aparece en ninguna forma, el contenido es insuficiente o existe pero es meramente superficial sin desarrollo real.

REGLAS PARA LA OBSERVACIÓN:
- Si PRESENTE: indica brevemente DÓNDE o CÓMO se evidencia (ej. "Se desarrolla en la diapositiva 3 con ejemplos prácticos").
- Si AUSENTE: indica QUÉ falta y UNA recomendación concreta para subsanarlo (ej. "No se menciona. Se recomienda agregar una sección explicando X con al menos un ejemplo").
- Sé específico y accionable. Evita respuestas genéricas como "no se encontró".

REQUISITOS A EVALUAR ({len(params)} en total):
{reqs_list}

CONTENIDO DEL DOCUMENTO:
{doc_snippet}

Responde ÚNICAMENTE con un array JSON válido (sin markdown ni texto extra), con exactamente {len(params)} objetos en el mismo orden que los requisitos:
[
  {{"presente": "Si" o "No", "observacion": "Observación específica y accionable de 1-2 oraciones", "confidence": 0.0-1.0}},
  ...
]"""

        # ── Prioridad 1: DeepSeek ─────────────────────────────────────────────
        if self._deepseek_key and getattr(self, '_req_use_deepseek', True):
            print("  🤖 Usando DeepSeek como proveedor principal...")
            ds_results = self._call_deepseek_batch(
                params, doc_snippet, section_name, group_name
            )
            if ds_results is not None:
                return ds_results
            print("  ⚠️  DeepSeek falló, intentando con Gemini...")
        elif self._deepseek_key and not getattr(self, '_req_use_deepseek', True):
            print("  ⏭️  DeepSeek desactivado por configuración")

        # ── Prioridad 2: Gemini con rotación de claves ────────────────────────
        available_keys = self._get_available_keys()
        target_model = model_name or GEMINI_MODEL

        # Distingue entre dos tipos de 429:
        #   A) Cuota diaria agotada (limit: 0) → rotar key inmediatamente
        #   B) Burst de RPM                    → reintentar misma key con backoff
        MAX_RPM_RETRIES = 4   # reintentos por burst: 5s, 10s, 20s, 40s

        for key in available_keys:
            short_key = key[:8] + '...'
            backoff = 5.0   # primer wait: 5s (estándar de la industria)

            for rpm_attempt in range(MAX_RPM_RETRIES + 1):
                try:
                    raw = self._call_gemini_raw(key, target_model, prompt)
                    parsed = json.loads(raw)

                    if not isinstance(parsed, list):
                        raise ValueError("Respuesta no es lista")

                    results = []
                    for i, p in enumerate(params):
                        if i < len(parsed):
                            item = parsed[i]
                            raw_conf = item.get('confidence', None)
                            try:
                                conf = round(float(raw_conf), 3) if raw_conf is not None else None
                            except (TypeError, ValueError):
                                conf = None
                            results.append({
                                'presente':    str(item.get('presente', 'No')),
                                'observacion': str(item.get('observacion', '')),
                                'confidence':  conf,
                            })
                        else:
                            results.append(self._fallback_keyword_check(p['sub_seccion'], doc_text))
                    key_idx = self._api_keys.index(key) + 1 if key in self._api_keys else '?'
                    print(f"  🤖 [MODELO: GEMINI {target_model} / Key #{key_idx}] respondió — {len(results)} resultados")
                    return results

                except json.JSONDecodeError:
                    print(f"  ⚠️  Batch JSON inválido ({short_key})")
                    break   # JSON raro → pasar al siguiente proveedor

                except Exception as e:
                    err_str = str(e)
                    is_quota = (
                        '429' in err_str
                        or 'quota' in err_str.lower()
                        or 'RESOURCE_EXHAUSTED' in err_str
                    )
                    if not is_quota:
                        print(f"  ⚠️  Error Gemini: {e}")
                        break

                    # — Distinguir: cuota diaria vs. burst de RPM —
                    daily_exhausted = 'limit: 0' in err_str or 'limit_per_day' in err_str.lower()

                    if daily_exhausted:
                        # Cuota diaria: esta key no sirve más hoy → rotar
                        self._exhausted_keys.add(key)
                        remaining = len(self._get_available_keys())
                        print(f"  🔄 Key ({short_key}) cuota diaria agotada "
                              f"→ rotando ({remaining} key{'s' if remaining != 1 else ''} restante{'s' if remaining != 1 else ''})")
                        break   # salir del bucle RPM, ir a siguiente key

                    else:
                        # Burst de RPM: la key sigue válida, solo hay que esperar
                        if rpm_attempt >= MAX_RPM_RETRIES:
                            # Demasiados reintentos de RPM → tratar como agotada
                            self._exhausted_keys.add(key)
                            print(f"  ⚠️  Key ({short_key}) no respondió tras {MAX_RPM_RETRIES} reintentos RPM → rotando")
                            break
                        wait = min(backoff, 60.0)   # máximo 60s de espera
                        print(f"  ⏳ RPM excedido ({short_key}) — esperando {wait:.0f}s "
                              f"(intento {rpm_attempt + 1}/{MAX_RPM_RETRIES})...")
                        time.sleep(wait)
                        backoff *= 2   # doblar: 5 → 10 → 20 → 40s

        # ── Respaldo 3: Groq ──────────────────────────────────────────────────
        if self._groq_key and getattr(self, '_req_use_groq', True):
            print("  🔄 Intentando con Groq (respaldo)...")
            groq_results = self._call_groq_batch(
                params, doc_snippet, section_name, group_name
            )
            if groq_results is not None:
                return groq_results
        elif self._groq_key and not getattr(self, '_req_use_groq', True):
            print("  ⏭️  Groq desactivado por configuración")

        # ── Respaldo 4: OpenRouter ────────────────────────────────────────────
        if self._openrouter_key and getattr(self, '_req_use_openrouter', True):
            print("  🔄 Intentando con OpenRouter (respaldo)...")
            or_results = self._call_openrouter_batch(
                params, doc_snippet, section_name, group_name
            )
            if or_results is not None:
                return or_results
        elif self._openrouter_key and not getattr(self, '_req_use_openrouter', True):
            print("  ⏭️  OpenRouter desactivado por configuración")

        # ── Último recurso: keyword matching ─────────────────────────────────
        return [self._fallback_keyword_check(p['sub_seccion'], doc_text) for p in params]

    # ──────────────────────────────────────────────────────────────────────────
    # Generación del reporte Excel (archivo nuevo, no modifica la matriz)
    # ──────────────────────────────────────────────────────────────────────────

    REPORT_PREFIX = "Reporte Validacion"

    def _build_report_excel(
        self,
        section_name: str,
        requirements: List[Dict],
        results: List[Dict],
        compliance_pct: float,
        docs_analyzed: List[str],
        model_name: str,
        groups: Optional[List[Dict]] = None,
    ) -> bytes:
        """
        Crea un nuevo workbook de reporte desde cero.
        No toca el Excel original de la matriz.
        """
        from datetime import datetime
        from openpyxl.styles import (
            Font, PatternFill, Alignment, Border, Side
        )

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Validación Contenido"

        # ── Estilos ──────────────────────────────────────────────
        title_font   = Font(bold=True, size=14)
        header_font  = Font(bold=True, size=10, color="FFFFFF")
        label_font   = Font(bold=True, size=10)
        normal_font  = Font(size=10)

        header_fill  = PatternFill("solid", fgColor="1E3A5F")   # azul oscuro
        green_fill   = PatternFill("solid", fgColor="C6EFCE")   # verde claro
        red_fill     = PatternFill("solid", fgColor="FFC7CE")   # rojo claro
        gray_fill    = PatternFill("solid", fgColor="F2F2F2")   # gris info

        center = Alignment(horizontal="center", vertical="center", wrap_text=True)
        left   = Alignment(horizontal="left",   vertical="center", wrap_text=True)

        thin = Side(style="thin", color="CCCCCC")
        border = Border(left=thin, right=thin, top=thin, bottom=thin)

        # ── Fila 1: Título ────────────────────────────────────────
        ws.merge_cells("A1:D1")
        ws["A1"].value     = "Reporte de Validación de Contenido"
        ws["A1"].font      = title_font
        ws["A1"].alignment = center
        ws.row_dimensions[1].height = 28

        # ── Filas 2-4: Metadatos ──────────────────────────────────
        meta = [
            ("Sección:",      section_name,
             "Fecha:",        datetime.now().strftime("%Y-%m-%d %H:%M")),
            ("Cumplimiento:", f"{compliance_pct}%",
             "Modelo IA:",    model_name),
            ("Documentos:",   ", ".join(docs_analyzed) if docs_analyzed else "—",
             "",              ""),
        ]
        for r_offset, (la, va, lb, vb) in enumerate(meta, start=2):
            ws[f"A{r_offset}"].value      = la
            ws[f"A{r_offset}"].font       = label_font
            ws[f"A{r_offset}"].fill       = gray_fill
            ws[f"B{r_offset}"].value      = va
            ws[f"B{r_offset}"].font       = normal_font
            ws[f"C{r_offset}"].value      = lb
            ws[f"C{r_offset}"].font       = label_font
            ws[f"C{r_offset}"].fill       = gray_fill
            ws[f"D{r_offset}"].value      = vb
            ws[f"D{r_offset}"].font       = normal_font
            ws.row_dimensions[r_offset].height = 18

        # Merge celda B3 si "Documentos" es larga
        ws.merge_cells("B4:D4")

        group_header_fill = PatternFill("solid", fgColor="2563EB")   # azul medio (encab. grupo)
        group_header_font = Font(bold=True, size=10, color="FFFFFF")

        # ── Fila 5: Espaciador ────────────────────────────────────
        ws.row_dimensions[5].height = 8

        # ── Fila 6: Encabezados de tabla ──────────────────────────
        headers = ["Sub-sección / Requisito", "Archivo analizado", "Presente", "Observación"]
        for col, h in enumerate(headers, start=1):
            cell = ws.cell(row=6, column=col)
            cell.value     = h
            cell.font      = header_font
            cell.fill      = header_fill
            cell.alignment = center
            cell.border    = border
        ws.row_dimensions[6].height = 20

        # ── Filas 7+: Datos (agrupados si hay grupos, planos si no) ──────────
        current_row = 7

        if groups:
            for group in groups:
                # Fila de encabezado de grupo (tipo de documento)
                g_pct   = group.get('compliance_percentage', 0)
                g_file  = group.get('matched_file') or '— archivo no encontrado —'
                g_label = (f"📑 {group['group_name']}  ·  {g_file}  "
                           f"·  {group.get('present_count', 0)}/{group.get('total_params', 0)} "
                           f"({g_pct}%)")
                ws.merge_cells(
                    start_row=current_row, start_column=1,
                    end_row=current_row,   end_column=4
                )
                gc = ws.cell(row=current_row, column=1)
                gc.value     = g_label
                gc.font      = group_header_font
                gc.fill      = group_header_fill
                gc.alignment = left
                ws.row_dimensions[current_row].height = 18
                current_row += 1

                # Filas de parámetros del grupo
                for r in group.get('results', []):
                    presente   = r.get('presente', 'No')
                    row_fill   = green_fill if presente == 'Si' else red_fill
                    row_data   = [r['sub_seccion'], g_file, presente, r.get('observacion', '')]
                    for col, val in enumerate(row_data, start=1):
                        cell = ws.cell(row=current_row, column=col)
                        cell.value     = val
                        cell.font      = normal_font
                        cell.fill      = row_fill
                        cell.alignment = left if col != 3 else center
                        cell.border    = border
                    ws.row_dimensions[current_row].height = 32
                    current_row += 1

        else:
            # Fallback: lista plana (sin grupos)
            for req, result in zip(requirements, results):
                presente = result.get('presente', 'No')
                row_fill = green_fill if presente == 'Si' else red_fill
                row_data = [req['sub_seccion'], req.get('autor', ''), presente, result.get('observacion', '')]
                for col, val in enumerate(row_data, start=1):
                    cell = ws.cell(row=current_row, column=col)
                    cell.value     = val
                    cell.font      = normal_font
                    cell.fill      = row_fill
                    cell.alignment = left if col != 3 else center
                    cell.border    = border
                ws.row_dimensions[current_row].height = 36
                current_row += 1

        # ── Anchos de columna ─────────────────────────────────────
        ws.column_dimensions["A"].width = 55
        ws.column_dimensions["B"].width = 20
        ws.column_dimensions["C"].width = 12
        ws.column_dimensions["D"].width = 60

        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    def _save_report(
        self,
        report_bytes: bytes,
        section_name: str,
        target_folder_id: str,
    ) -> Optional[Dict]:
        """
        Sube el reporte a Drive.
        Si ya existe un reporte para esta sección en la carpeta, lo actualiza.
        Si no, crea uno nuevo.
        Retorna {'id', 'name', 'webViewLink'} o None.
        """
        from datetime import datetime

        # Prefijo estable para poder encontrarlo después
        safe_section = section_name.replace(" ", "_")
        prefix = f"{self.REPORT_PREFIX} {safe_section}"

        # Buscar si ya existe
        existing = drive_service.find_file_by_prefix(target_folder_id, prefix)

        if existing:
            # Actualizar en lugar de crear nuevo
            ok = drive_service.upload_file(report_bytes, EXCEL_MIME_TYPE, existing['id'])
            if ok:
                print(f"  📤 Reporte actualizado: '{existing['name']}'")
                return existing
            return None
        else:
            # Crear archivo nuevo con timestamp
            ts = datetime.now().strftime("%Y-%m-%d %H-%M")
            filename = f"{prefix} - {ts}.xlsx"
            result = drive_service.create_file(
                report_bytes, EXCEL_MIME_TYPE, filename, target_folder_id
            )
            return result

    def _write_results_to_matrix(
        self,
        ws,
        col_map: Dict[str, int],
        group_results: List[Dict],
    ) -> int:
        """
        Escribe los resultados de validación directamente en la hoja 'Matriz observaciones'.
        Actualiza las columnas 'Presente' y 'Observaciones' para cada parámetro usando
        el row_idx guardado durante el parseo.

        Aplica fill verde (C6EFCE) para 'Si' y rojo (FFC7CE) para 'No'.
        Retorna el número de filas actualizadas.
        """
        from openpyxl.styles import PatternFill, Alignment

        presente_col = col_map.get('presente')
        obs_col      = col_map.get('observaciones')
        autor_col    = col_map.get('autor')

        if not presente_col and not obs_col:
            print("  ⚠️  No se encontraron columnas 'Presente'/'Observaciones' — no se puede escribir en la matriz")
            return 0

        green_fill = PatternFill("solid", fgColor="C6EFCE")
        red_fill   = PatternFill("solid", fgColor="FFC7CE")
        left_wrap  = Alignment(horizontal="left", vertical="center", wrap_text=True)
        center     = Alignment(horizontal="center", vertical="center")

        print(f"  🖊️  Escribiendo {len(group_results)} grupo(s) | "
              f"col Presente={presente_col} | col Obs={obs_col} | col Autor={autor_col}")
        updated = 0
        for group in group_results:
            group_rows = [r.get('row_idx') for r in group.get('results', []) if r.get('row_idx')]
            print(f"     ✏️  '{group.get('group_name','?')}': "
                  f"{len(group.get('results',[]))} resultado(s) | "
                  f"filas={group_rows[:5]}{'…' if len(group_rows)>5 else ''}")

            # Escribir autor en la fila del encabezado del grupo
            if autor_col and group.get('header_row') and group.get('autor'):
                autor_cell = ws.cell(row=group['header_row'], column=autor_col)
                autor_cell.value = group['autor']
                autor_cell.alignment = center
                print(f"        👤 Autor '{group['autor']}' → fila {group['header_row']}, col {autor_col}")

            for r in group.get('results', []):
                row_idx = r.get('row_idx')
                if not row_idx:
                    print(f"        ⚠️  row_idx ausente en resultado: {r.get('sub_seccion','?')[:40]}")
                    continue

                presente = r.get('presente', 'No')
                fill     = green_fill if presente == 'Si' else red_fill

                if presente_col:
                    cell       = ws.cell(row=row_idx, column=presente_col)
                    cell.value = presente
                    cell.fill  = fill
                    cell.alignment = center

                if obs_col:
                    obs_cell       = ws.cell(row=row_idx, column=obs_col)
                    obs_cell.value = r.get('observacion', '')
                    obs_cell.fill  = fill
                    obs_cell.alignment = left_wrap

                updated += 1

        return updated

    def _evaluate_lab_group_batch(
        self,
        params: List[Dict],
        doc_text: str,
        section_name: str,
        group_name: str,
        folder_type: str,
        use_gemini: bool = True,
        model_name: Optional[str] = None,
    ) -> Tuple[List[Dict[str, str]], Dict[str, str]]:
        """
        Evaluación enriquecida para carpetas Proyectos/Practicas/Tareas.
        Retorna (requirements_results, document_analysis).

        requirements_results: [{presente, observacion}, ...]
        document_analysis: {descripcion, enfoque, complejidad}
        """
        lab_info = LAB_FOLDER_DESCRIPTIONS.get(folder_type, {})
        tipo_doc  = lab_info.get('tipo', group_name)
        max_chars = 15_000
        doc_snippet = doc_text[:max_chars]

        reqs_list = "\n".join(
            f"{i+1}. {p['sub_seccion']}" for i, p in enumerate(params)
        )

        prompt = f"""Eres un evaluador académico universitario experto en revisión de trabajos estudiantiles. Analiza el documento y realiza DOS tareas:

═══════════════════════════════════════════════════
TAREA 1 — EVALUAR REQUISITOS DE LA RÚBRICA
═══════════════════════════════════════════════════
Determina si cada requisito está cubierto en el documento.

CRITERIOS:
• PRESENTE (Si): el documento aborda el requisito de forma explícita (nombrado directamente)
  O implícita (demostrado mediante ejemplos, código, diagramas, resultados o desarrollo
  práctico que evidencia el concepto aunque use terminología distinta).
• AUSENTE (No): el requisito no aparece, o aparece de forma tan superficial que no
  aporta valor real al trabajo.

INSTRUCCIONES PARA CADA REQUISITO:
- "presente": "Si" o "No" (según los criterios anteriores).
- "observacion":
    Si PRESENTE → cita específicamente dónde o cómo se evidencia en el documento
    (ej: "Se presenta en la sección Metodología con el diagrama de flujo del proceso X").
    Si AUSENTE → explica qué falta y por qué es importante para este tipo de {tipo_doc}.
- "sugerencia": SOLO si AUSENTE → indica una acción concreta y específica para
    incluir este elemento (ej: "Agregar una sección de Conclusiones que relacione
    los resultados obtenidos con los objetivos planteados").
- "confidence": 0.0–1.0 (tu certeza; >0.85 solo si la evidencia es clara y directa).

═══════════════════════════════════════════════════
TAREA 2 — ANÁLISIS DEL DOCUMENTO
═══════════════════════════════════════════════════
Proporciona un análisis objetivo del trabajo:
- descripcion: qué desarrolla o presenta el documento, qué problema aborda (2-3 oraciones).
- enfoque: enfoque técnico/metodológico principal que usa el estudiante (1-2 oraciones).
- complejidad: nivel del trabajo ("Básica", "Intermedia" o "Avanzada") con una
  justificación breve basada en la profundidad, originalidad y cantidad de contenido.

═══════════════════════════════════════════════════
CONTEXTO DE EVALUACIÓN
═══════════════════════════════════════════════════
- Tipo de actividad: {tipo_doc}
- Carpeta / sección del curso: {section_name}
- Documento evaluado: {group_name}

REQUISITOS A EVALUAR ({len(params)} en total):
{reqs_list}

CONTENIDO DEL DOCUMENTO:
{doc_snippet}

═══════════════════════════════════════════════════
Responde ÚNICAMENTE con JSON válido (sin markdown ni texto extra):
{{
  "requirements": [
    {{"presente": "Si", "observacion": "Evidencia específica de cómo aparece", "confidence": 0.90}},
    {{"presente": "No", "observacion": "Qué falta y por qué importa", "sugerencia": "Acción concreta para incluirlo", "confidence": 0.85}},
    ... ({len(params)} elementos en el mismo orden que la lista)
  ],
  "document_analysis": {{
    "descripcion": "Descripción objetiva del trabajo (2-3 oraciones)",
    "enfoque": "Enfoque técnico/metodológico principal (1-2 oraciones)",
    "complejidad": "Básica|Intermedia|Avanzada — justificación en 1 oración"
  }}
}}"""

        empty_analysis = {'descripcion': '', 'enfoque': '', 'complejidad': ''}

        def _parse_lab_response(raw: str) -> Tuple[Optional[List], Optional[Dict]]:
            raw = re.sub(r"^```json?\s*", "", raw.strip())
            raw = re.sub(r"\s*```$", "", raw)
            try:
                data = json.loads(raw)
                reqs = data.get('requirements', [])
                analysis = data.get('document_analysis', {})
                if isinstance(reqs, list) and len(reqs) == len(params):
                    return reqs, analysis
            except Exception:
                pass
            return None, None

        def _pack(reqs: List, analysis: Optional[Dict]) -> Tuple[List[Dict], Dict]:
            packed = []
            for r in reqs:
                raw_conf = r.get('confidence', None)
                try:
                    conf = round(float(raw_conf), 3) if raw_conf is not None else None
                except (TypeError, ValueError):
                    conf = None
                presente = str(r.get('presente', 'No'))
                packed.append({
                    'presente':    presente,
                    'observacion': str(r.get('observacion', '')),
                    'sugerencia':  str(r.get('sugerencia', '')) if presente != 'Si' else '',
                    'confidence':  conf,
                })
            return packed, analysis or empty_analysis

        import urllib.request, urllib.error

        # ── Prioridad 1: DeepSeek ─────────────────────────────────────────────
        if self._deepseek_key:
            payload = json.dumps({
                "model": "deepseek-chat",
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.1,
                "max_tokens": 4000,
            }).encode("utf-8")
            try:
                req = urllib.request.Request(
                    "https://api.deepseek.com/v1/chat/completions",
                    data=payload,
                    headers={"Authorization": f"Bearer {self._deepseek_key}", "Content-Type": "application/json"},
                    method="POST",
                )
                with urllib.request.urlopen(req, timeout=90) as resp:
                    content = json.loads(resp.read().decode("utf-8"))["choices"][0]["message"]["content"].strip()
                reqs, analysis = _parse_lab_response(content)
                if reqs:
                    print(f"  🤖 [DeepSeek lab] respondió — complejidad: {(analysis or {}).get('complejidad','?')}")
                    return _pack(reqs, analysis)
                print("  ⚠️  DeepSeek lab: respuesta vacía o formato inesperado")
            except Exception as e:
                print(f"  ⚠️  DeepSeek lab: {e}")

        # ── Prioridad 2: Gemini con rotación de keys y retry RPM ──────────────
        MAX_RPM_RETRIES = 4
        target_model = model_name or GEMINI_MODEL

        for key in self._get_available_keys():
            short_key = key[:8] + '...'
            backoff = 5.0

            for rpm_attempt in range(MAX_RPM_RETRIES + 1):
                try:
                    raw = self._call_gemini_raw(key, target_model, prompt)
                    reqs, analysis = _parse_lab_response(raw)
                    if reqs:
                        key_idx = self._api_keys.index(key) + 1 if key in self._api_keys else '?'
                        print(f"  🤖 [Gemini lab / Key #{key_idx}] respondió — complejidad: {(analysis or {}).get('complejidad','?')}")
                        return _pack(reqs, analysis)
                    # Respuesta malformada → no reintentamos esta key
                    break

                except Exception as e:
                    err_str = str(e)
                    if '429' in err_str or 'quota' in err_str.lower():
                        # Distinguir cuota agotada (limit=0) de burst RPM
                        if 'limit' in err_str.lower() and ('0' in err_str or 'exhausted' in err_str.lower()):
                            self._exhausted_keys.add(key)
                            remaining = len([k for k in self._api_keys if k not in self._exhausted_keys])
                            print(f"  ⚠️  Cuota Gemini agotada ({short_key}) → rotando ({remaining} key(s) restante(s))")
                            break  # siguiente key
                        else:
                            if rpm_attempt >= MAX_RPM_RETRIES:
                                self._exhausted_keys.add(key)
                                print(f"  ⚠️  Key ({short_key}) sin respuesta tras {MAX_RPM_RETRIES} reintentos RPM → rotando")
                                break
                            wait = min(backoff, 60.0)
                            print(f"  ⏳ RPM excedido ({short_key}) — esperando {wait:.0f}s (intento {rpm_attempt + 1}/{MAX_RPM_RETRIES})...")
                            time.sleep(wait)
                            backoff *= 2
                    else:
                        print(f"  ⚠️  Gemini lab ({short_key}): {e}")
                        break  # error no RPM → siguiente key

        # ── Prioridad 3: Groq ─────────────────────────────────────────────────
        if self._groq_key:
            print("  🔄 [lab] Intentando con Groq (respaldo)...")
            groq_results = self._call_groq_batch(params, doc_snippet, section_name, group_name)
            if groq_results is not None:
                print(f"  🤖 [Groq lab] respondió — {len(groq_results)} resultados")
                return groq_results, empty_analysis

        # ── Prioridad 4: OpenRouter ───────────────────────────────────────────
        if self._openrouter_key:
            print("  🔄 [lab] Intentando con OpenRouter (respaldo)...")
            or_results = self._call_openrouter_batch(params, doc_snippet, section_name, group_name)
            if or_results is not None:
                print(f"  🤖 [OpenRouter lab] respondió — {len(or_results)} resultados")
                return or_results, empty_analysis

        # ── Último recurso: keyword matching ──────────────────────────────────
        print("  ⚠️  [lab] Todos los proveedores fallaron — usando keywords")
        fallback = [self._fallback_keyword_check(p['sub_seccion'], doc_text) for p in params]
        return fallback, empty_analysis

    # ──────────────────────────────────────────────────────────────────────────
    # Punto de entrada principal
    # ──────────────────────────────────────────────────────────────────────────

    def validate_folder_content(
        self,
        semana_folder_id: str,
        semana_folder_name: str,
        candidate_folder_ids: List[str],
        db=None,
        user_id=None,
    ) -> Dict[str, Any]:
        """
        Pipeline completo de validación de contenido para una carpeta Semana_X.

        Nuevo enfoque por grupos:
          - El Excel agrupa parámetros bajo encabezados de tipo de documento
            (Presentación, Lectura, Video…)
          - Cada tipo de documento se localiza en la carpeta Drive
          - Los parámetros se evalúan SOLO contra el archivo correspondiente
            → Gemini recibe contexto puro, sin mezclar documentos
        """
        try:
            print(f"\n🧠 Iniciando validación de contenido: '{semana_folder_name}'")

            # 0. Detectar tipo de carpeta
            folder_type = self._get_folder_type(semana_folder_name)
            is_lab = folder_type in LAB_FOLDER_TYPES
            print(f"  🗂️  Tipo de carpeta: {folder_type}")

            # 0b. Leer settings dinámicos desde BD (si db disponible)
            use_gemini      = self.enabled
            active_model_name = GEMINI_MODEL
            # Defaults por si no hay BD
            self._req_use_deepseek    = True
            self._req_use_groq        = True
            self._req_use_openrouter  = True
            self._req_temperature     = 0.05
            self._req_max_tokens      = 2000
            if db is not None:
                try:
                    from app.services.settings_service import settings_service
                    active_model_name = settings_service.get("gemini_model", db) or GEMINI_MODEL
                    self._req_use_deepseek   = settings_service.get_bool("deepseek_enabled", db)
                    self._req_use_groq       = settings_service.get_bool("groq_enabled", db)
                    self._req_use_openrouter = settings_service.get_bool("openrouter_enabled", db)
                    try:
                        self._req_temperature = float(settings_service.get("ai_temperature", db) or 0.05)
                    except (ValueError, TypeError):
                        pass
                    self._req_max_tokens = settings_service.get_int("ai_max_tokens", db, default=2000)
                    print(
                        f"  ⚙️  Proveedores: DeepSeek={'✓' if self._req_use_deepseek else '✗'} "
                        f"Gemini={'✓' if use_gemini else '✗'} (modelo: {active_model_name}) "
                        f"Groq={'✓' if self._req_use_groq else '✗'} "
                        f"OpenRouter={'✓' if self._req_use_openrouter else '✗'} "
                        f"| temp={self._req_temperature} max_tokens={self._req_max_tokens}"
                    )
                except Exception:
                    pass

            # 0c. Recargar API keys desde BD (prioriza keys personales del usuario)
            self._reload_keys_from_db(db, user_id=user_id)

                # Gemini queda habilitado si existe al menos una key disponible y
                # la configuración del sistema no lo ha desactivado.
            if db is not None:
                use_gemini = bool(self._api_keys) and settings_service.get_bool("gemini_enabled", db)
            else:
                use_gemini = bool(self._api_keys)

            # 1. Derivar sección
            section_name = self._derive_section_from_folder(semana_folder_name)
            print(f"  📌 Sección: '{section_name}'")

            # 2. Buscar el Excel en cada carpeta candidata (busca en '0. Revision de Material' primero)
            matrix_info = None
            matrix_folder_id = None
            for fid in candidate_folder_ids:
                print(f"  🔎 Buscando matriz en: {fid}")
                matrix_info = self.find_matrix_file(fid)
                if matrix_info:
                    matrix_folder_id = fid
                    break

            if not matrix_info:
                return {
                    'success': False,
                    'error': f'No se encontró "{MATRIX_FILE_PREFIX}" en ninguna carpeta candidata'
                }
            matrix_file_id, matrix_file_name = matrix_info
            print(f"  📊 Matriz: '{matrix_file_name}'")

            # 3. Cargar workbook
            wb = self.read_workbook(matrix_file_id)
            if wb is None:
                return {'success': False, 'error': 'No se pudo leer el Excel'}

            # 3b. Seleccionar hoja y parsear requisitos según tipo de carpeta
            groups: List[Dict] = []
            col_map: Dict[str, int] = {}

            if is_lab:
                # ── Hojas Proyectos / Practicas / Tareas ─────────────────────
                ws = self.get_lab_sheet(wb, folder_type)
                if ws is None:
                    return {
                        'success': False,
                        'error': f'No se encontró la hoja "{LAB_FOLDER_TYPES[folder_type]}" en la Matriz. '
                                 f'Hojas disponibles: {wb.sheetnames}'
                    }
                # Parser dedicado: detecta columnas y agrupa en un solo paso
                groups, col_map = self._parse_lab_sheet_full(ws)

                if not groups:
                    return {
                        'success': False,
                        'error': (f'No se encontraron criterios en la hoja '
                                  f'"{LAB_FOLDER_TYPES[folder_type]}". '
                                  f'Verifica que la hoja tenga encabezados de grupo '
                                  f'(ej. "Proyecto 1") y criterios de evaluación.')
                    }
            else:
                # ── Hoja Semana (flujo original) ──────────────────────────────
                ws = self.get_content_sheet(wb)
                if ws is None:
                    return {'success': False, 'error': 'No se encontró la hoja "Matriz observaciones"'}

                try:
                    header_row_idx, col_map = self._find_header_row(ws)
                except ValueError as hdr_err:
                    print(f"  ⚠️  No se detectaron encabezados estándar: {hdr_err}")
                    return {'success': False, 'error': f'Error leyendo encabezados del Excel: {hdr_err}'}

                # 4. Parsear requisitos agrupados por tipo de documento
                groups = self.parse_requirements_by_group(ws, section_name, col_map, header_row_idx)
                if not groups:
                    return {
                        'success': False,
                        'error': f'No se encontraron requisitos para la sección "{section_name}"'
                    }

            # 5. Descargar y extraer texto de TODOS los documentos de la carpeta
            doc_texts, files_metadata = self.download_and_extract_all_docs(semana_folder_id)
            if not doc_texts:
                return {
                    'success': False,
                    'error': 'No se encontraron documentos PDF/DOCX/PPTX en la carpeta'
                }
            print(f"  📁 {len(doc_texts)} documento(s) en la carpeta")

            # 6. Validar cada grupo contra su documento específico
            group_results: List[Dict] = []
            already_matched: set = set()      # evitar asignar el mismo archivo a 2 grupos
            total_present = 0
            total_absent  = 0

            for group in groups:
                # 6a. Buscar el archivo correspondiente al tipo de documento.
                # Caso especial: grupo "General" (viene del raw fallback) →
                # concatenar TODOS los documentos de la carpeta.
                if self._normalize(group['group_name']) == 'general':
                    doc_text   = '\n\n---\n\n'.join(doc_texts.values())
                    file_found = len(doc_texts) > 0
                    matched_name = f"{len(doc_texts)} documento(s)"
                else:
                    matched_file = self._match_file_to_group(
                        group['group_name'], files_metadata, already_matched
                    )
                    file_found   = matched_file is not None
                    matched_name = matched_file['name'] if matched_file else None
                    if matched_name:
                        already_matched.add(matched_name)
                    doc_text = doc_texts.get(matched_name, '') if matched_name else ''

                print(f"\n  📑 Grupo '{group['group_name']}' → "
                      f"{'archivo: ' + matched_name if file_found else '⚠️ sin archivo'}")

                # 6b. Evaluar TODOS los parámetros del grupo en UNA sola llamada batch
                params_results: List[Dict] = []
                present_in_group = 0
                absent_in_group  = 0
                n = len(group['params'])

                doc_analysis: Dict[str, str] = {}

                if not file_found or not doc_text.strip():
                    # Sin archivo → todos ausentes
                    batch_evals = [
                        {
                            'presente':    'No',
                            'observacion': (f'Documento tipo "{group["group_name"]}" no encontrado '
                                           f'en la carpeta "{semana_folder_name}"'),
                        }
                        for _ in group['params']
                    ]
                elif is_lab:
                    print(f"    🔬 Batch lab ({folder_type}): {n} parámetros + análisis documental...")
                    batch_evals, doc_analysis = self._evaluate_lab_group_batch(
                        params       = group['params'],
                        doc_text     = doc_text,
                        section_name = section_name,
                        group_name   = group['group_name'],
                        folder_type  = folder_type,
                        use_gemini   = use_gemini,
                        model_name   = active_model_name,
                    )
                    time.sleep(4.5)
                else:
                    print(f"    🤖 Batch Gemini: {n} parámetros en 1 llamada...")
                    batch_evals = self._evaluate_group_batch(
                        params       = group['params'],
                        doc_text     = doc_text,
                        section_name = section_name,
                        group_name   = group['group_name'],
                        use_gemini   = use_gemini,
                        model_name   = active_model_name,
                    )
                    time.sleep(4.5)

                for i, (param, eval_r) in enumerate(zip(group['params'], batch_evals), 1):
                    sub = param['sub_seccion']
                    print(f"    [{i}/{n}] '{sub[:50]}' → {eval_r['presente']}")
                    if eval_r['presente'] == 'Si':
                        present_in_group += 1
                    else:
                        absent_in_group += 1
                    params_results.append({
                        'sub_seccion': sub,
                        'autor':       param['autor'],
                        'presente':    eval_r['presente'],
                        'observacion': eval_r['observacion'],
                        'row_idx':     param['row_idx'],
                    })

                total_present += present_in_group
                total_absent  += absent_in_group
                g_total = len(group['params'])
                g_pct   = round(present_in_group / g_total * 100, 2) if g_total else 0

                owner_name = ''
                if self._normalize(group['group_name']) != 'general' and matched_file:
                    owner_name = matched_file.get('owner', '')

                group_entry = {
                    'group_name':            group['group_name'],
                    'matched_file':          matched_name,
                    'file_found':            file_found,
                    'total_params':          g_total,
                    'present_count':         present_in_group,
                    'absent_count':          absent_in_group,
                    'compliance_percentage': g_pct,
                    'results':               params_results,
                    'autor':                 owner_name,
                    'header_row':            group.get('header_row'),
                }
                # Campos extra para carpetas lab
                if is_lab and doc_analysis:
                    group_entry['descripcion'] = doc_analysis.get('descripcion', '')
                    group_entry['enfoque']     = doc_analysis.get('enfoque', '')
                    group_entry['complejidad'] = doc_analysis.get('complejidad', '')

                group_results.append(group_entry)
                print(f"  ✅ '{group['group_name']}': {present_in_group}/{g_total} ({g_pct}%)"
                      + (f" | complejidad: {doc_analysis.get('complejidad','')}" if is_lab and doc_analysis else ""))

            # 7. Estadísticas globales
            total_requirements = total_present + total_absent
            compliance_pct = round(total_present / total_requirements * 100, 2) if total_requirements else 0

            # Lista plana para compatibilidad (report, batch, etc.)
            flat_results = [r for g in group_results for r in g['results']]
            # Flat requirements list for report builder
            flat_requirements = [
                {'sub_seccion': r['sub_seccion'], 'autor': r['autor']}
                for r in flat_results
            ]

            print(f"\n✅ Validación: {total_present}/{total_requirements} ({compliance_pct}%) presentes "
                  f"en {len(groups)} grupo(s)")

            # 7b. Escribir resultados en la MATRIZ ORIGINAL (Presente + Observaciones)
            print("\n  📝 Escribiendo resultados en la matriz original...")
            n_written = self._write_results_to_matrix(ws, col_map, group_results)
            matrix_updated = False
            if n_written > 0:
                try:
                    buf = io.BytesIO()
                    wb.save(buf)
                    matrix_bytes = buf.getvalue()
                    matrix_updated = drive_service.upload_file(
                        matrix_bytes, EXCEL_MIME_TYPE, matrix_file_id
                    )
                    if matrix_updated:
                        print(f"  ✅ Matriz actualizada en Drive: {n_written} fila(s) escritas")
                    else:
                        print("  ⚠️  No se pudo subir la matriz actualizada a Drive")
                except Exception as e:
                    print(f"  ⚠️  Error guardando matriz: {e}")

            # 8. Generar reporte Excel (sin tocar la matriz original)
            report_bytes = self._build_report_excel(
                section_name   = section_name,
                requirements   = flat_requirements,
                results        = flat_results,
                compliance_pct = compliance_pct,
                docs_analyzed  = list(doc_texts.keys()),
                model_name     = active_model_name,
                groups         = group_results,
            )
            report_info = self._save_report(report_bytes, section_name, semana_folder_id)
            if report_info:
                print(f"  📄 Reporte: '{report_info.get('name')}'")

            return {
                'success':               True,
                'section':               section_name,
                'folder_type':           folder_type,
                'is_lab':                is_lab,
                'total_requirements':    total_requirements,
                'present_count':         total_present,
                'absent_count':          total_absent,
                'compliance_percentage': compliance_pct,
                'groups':                group_results,
                'results':               flat_results,
                'documents_analyzed':    list(doc_texts.keys()),
                'report_generated':      report_info is not None,
                'report_name':           report_info.get('name') if report_info else None,
                'report_link':           report_info.get('webViewLink') if report_info else None,
                'excel_updated':         matrix_updated,
                'gemini_enabled':        self.enabled,
            }

        except Exception as e:
            print(f"❌ Error en validación de contenido: {e}")
            import traceback
            traceback.print_exc()
            return {'success': False, 'error': str(e)}


# Singleton global
document_content_validation_service = DocumentContentValidationService()
